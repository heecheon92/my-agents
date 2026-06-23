"""Local Office upload parsing into Markdown plus source-location metadata."""

from __future__ import annotations

import hashlib
import logging
import re
import zipfile
from collections.abc import Sequence
from dataclasses import dataclass, field
from importlib.metadata import PackageNotFoundError, version
from io import BytesIO
from typing import Any

from openpyxl import load_workbook
from openpyxl.utils import get_column_letter
from pptx import Presentation

from my_agents.knowledge.pdf_uploads import (
    MAX_PDF_UPLOAD_BYTES,
    DoclingExtractionConfig,
    _docling_accelerator_device,
)

MAX_OFFICE_UPLOAD_BYTES = MAX_PDF_UPLOAD_BYTES
MAX_OFFICE_ARCHIVE_MEMBERS = 600
MAX_OFFICE_ARCHIVE_UNCOMPRESSED_BYTES = 25 * 1024 * 1024
MAX_OFFICE_ARCHIVE_COMPRESSION_RATIO = 100
MAX_EXCEL_WORKSHEETS = 50
MAX_EXCEL_ROWS_PER_SHEET = 5000
MAX_EXCEL_NON_EMPTY_CELLS = 50000
MAX_EXCEL_ROW_RECORDS_PER_SHEET = 1000
MAX_POWERPOINT_SLIDES = 200
MAX_POWERPOINT_SHAPES = 2000
MAX_POWERPOINT_TABLE_CELLS = 20000
MAX_OFFICE_MARKDOWN_CHARS = 500_000
EXCEL_UPLOAD_PARSER_NAME = "openpyxl_markdown_v1"
POWERPOINT_UPLOAD_PARSER_NAME = "python_pptx_markdown_v1"
WORD_DOCX_UPLOAD_PARSER_NAME = "docling_docx_markdown_v1"
EXCEL_SOURCE_TYPE = "spreadsheet"
POWERPOINT_SOURCE_TYPE = "presentation"
WORD_DOCUMENT_SOURCE_TYPE = "word_document"
EXCEL_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
POWERPOINT_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
WORD_DOCX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
OFFICE_PARSER_PROVIDER = "local"
OFFICE_PARSER_MODE = "deterministic"

_EXCEL_SUFFIX = ".xlsx"
_POWERPOINT_SUFFIX = ".pptx"
_WORD_DOCX_SUFFIX = ".docx"
_SUPPORTED_OFFICE_SUFFIXES = frozenset({_EXCEL_SUFFIX, _POWERPOINT_SUFFIX, _WORD_DOCX_SUFFIX})
_SUPPORTED_OFFICE_SUFFIX_LIST = ".xlsx, .pptx, and .docx"
_GENERIC_UPLOAD_CONTENT_TYPES = frozenset({"", "application/octet-stream"})
_EXCEL_CONTENT_TYPES = frozenset({EXCEL_CONTENT_TYPE, *_GENERIC_UPLOAD_CONTENT_TYPES})
_POWERPOINT_CONTENT_TYPES = frozenset({POWERPOINT_CONTENT_TYPE, *_GENERIC_UPLOAD_CONTENT_TYPES})
_WORD_DOCX_CONTENT_TYPES = frozenset({WORD_DOCX_CONTENT_TYPE, *_GENERIC_UPLOAD_CONTENT_TYPES})
_MONTH_VALUE_PATTERN = re.compile(r"^(?:[1-9]|1[0-2])월$")
_NUMERIC_VALUE_PATTERN = re.compile(r"^[+-]?(?:(?:\d{1,3}(?:,\d{3})+)|\d+)(?:\.\d+)?%?$")
_MARKDOWN_HEADING_PATTERN = re.compile(r"^(?P<marks>#{1,6})\s+(?P<title>.+?)\s*$")
_MARKDOWN_ORDERED_LIST_PATTERN = re.compile(r"^\s*\d+[\.)]\s+")
logger = logging.getLogger(__name__)


class OfficeUploadError(ValueError):
    """Raised when a local Office upload cannot be safely parsed."""


@dataclass(frozen=True)
class ParsedOfficeDocument:
    """Markdown and parser artifact metadata extracted from an Office upload."""

    content: str
    source_type: str
    source_content_type: str
    byte_size: int
    sha256: str
    parser_name: str
    parser_provider: str = OFFICE_PARSER_PROVIDER
    parser_version: str | None = None
    parser_mode: str = OFFICE_PARSER_MODE
    elements: list[dict[str, Any]] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)


@dataclass
class _MarkdownBuilder:
    parts: list[str] = field(default_factory=list)
    cursor: int = 0

    def append(self, text: str) -> tuple[int, int]:
        start = self.cursor
        self.parts.append(text)
        self.cursor += len(text)
        return start, self.cursor

    def content(self) -> str:
        return "".join(self.parts).strip()


@dataclass(frozen=True)
class _WorksheetGrid:
    rows: list[tuple[int, list[str]]]
    first_column_index: int
    last_column_index: int
    merged_cell_count: int


@dataclass(frozen=True)
class _WorksheetMarkdown:
    markdown: str
    elements: list[dict[str, Any]]
    warnings: list[str] = field(default_factory=list)


def parse_uploaded_office_document(
    *,
    filename: str | None,
    content_type: str | None,
    content: bytes,
    docling_config: DoclingExtractionConfig | None = None,
) -> ParsedOfficeDocument:
    """Parse a supported OOXML Office upload into canonical Markdown and metadata."""
    safe_filename = _validate_office_filename(filename)
    suffix = _filename_suffix(safe_filename)
    logger.info(
        "office_upload.dispatch filename=%s suffix=%s content_type=%s bytes=%d sha256=%s",
        safe_filename,
        suffix,
        content_type,
        len(content),
        hashlib.sha256(content).hexdigest(),
    )
    if suffix == _EXCEL_SUFFIX:
        return parse_uploaded_xlsx(
            filename=safe_filename, content_type=content_type, content=content
        )
    if suffix == _POWERPOINT_SUFFIX:
        return parse_uploaded_pptx(
            filename=safe_filename, content_type=content_type, content=content
        )
    if suffix == _WORD_DOCX_SUFFIX:
        return parse_uploaded_docx(
            filename=safe_filename,
            content_type=content_type,
            content=content,
            docling_config=docling_config,
        )
    raise OfficeUploadError(f"only {_SUPPORTED_OFFICE_SUFFIX_LIST} Office uploads are supported")


def parse_uploaded_xlsx(
    *,
    filename: str | None,
    content_type: str | None,
    content: bytes,
) -> ParsedOfficeDocument:
    """Parse a local `.xlsx` workbook into Markdown tables and sheet locations."""
    safe_filename = _validate_office_filename(filename, expected_suffix=_EXCEL_SUFFIX)
    _validate_content_type(
        suffix=_EXCEL_SUFFIX,
        content_type=content_type,
        allowed_content_types=_EXCEL_CONTENT_TYPES,
        expected_content_type=EXCEL_CONTENT_TYPE,
    )
    _validate_office_content(content, filename=safe_filename)
    _validate_ooxml_archive(content, filename=safe_filename)
    try:
        workbook = load_workbook(BytesIO(content), data_only=True, read_only=False)
    except Exception as exc:
        raise OfficeUploadError(
            f"{safe_filename} could not be parsed as a valid .xlsx file"
        ) from exc

    builder = _MarkdownBuilder()
    elements: list[dict[str, Any]] = []
    warnings: list[str] = []
    if len(workbook.worksheets) > MAX_EXCEL_WORKSHEETS:
        raise OfficeUploadError("uploaded spreadsheet exceeds the sheet count limit")
    for worksheet in workbook.worksheets:
        section = _worksheet_markdown(worksheet)
        if section is None:
            warnings.append(f"worksheet {worksheet.title!r} has no extractable cells")
            continue
        if builder.cursor:
            builder.append("\n\n")
        section_start, _ = builder.append(section.markdown)
        _validate_markdown_budget(builder)
        warnings.extend(section.warnings)
        for element in section.elements:
            elements.append(
                _source_element(
                    kind=element["kind"],
                    markdown_start=section_start + element["markdown_start"],
                    markdown_end=section_start + element["markdown_end"],
                    source_location=element["source_location"],
                )
            )

    markdown_content = builder.content()
    if not markdown_content:
        raise OfficeUploadError("uploaded spreadsheet has no extractable cells")
    return ParsedOfficeDocument(
        content=markdown_content,
        source_type=EXCEL_SOURCE_TYPE,
        source_content_type=EXCEL_CONTENT_TYPE,
        byte_size=len(content),
        sha256=hashlib.sha256(content).hexdigest(),
        parser_name=EXCEL_UPLOAD_PARSER_NAME,
        parser_version=_package_version("openpyxl"),
        elements=elements,
        warnings=warnings,
    )


def parse_uploaded_pptx(
    *,
    filename: str | None,
    content_type: str | None,
    content: bytes,
) -> ParsedOfficeDocument:
    """Parse a local `.pptx` deck into Markdown sections and slide locations."""
    safe_filename = _validate_office_filename(filename, expected_suffix=_POWERPOINT_SUFFIX)
    _validate_content_type(
        suffix=_POWERPOINT_SUFFIX,
        content_type=content_type,
        allowed_content_types=_POWERPOINT_CONTENT_TYPES,
        expected_content_type=POWERPOINT_CONTENT_TYPE,
    )
    _validate_office_content(content, filename=safe_filename)
    _validate_ooxml_archive(content, filename=safe_filename)
    try:
        presentation = Presentation(BytesIO(content))
    except Exception as exc:
        raise OfficeUploadError(
            f"{safe_filename} could not be parsed as a valid .pptx file"
        ) from exc

    builder = _MarkdownBuilder()
    elements: list[dict[str, Any]] = []
    warnings: list[str] = []
    if len(presentation.slides) > MAX_POWERPOINT_SLIDES:
        raise OfficeUploadError("uploaded presentation exceeds the slide count limit")
    shape_count = 0
    table_cell_count = 0
    for slide_number, slide in enumerate(presentation.slides, start=1):
        shape_count += len(slide.shapes)
        if shape_count > MAX_POWERPOINT_SHAPES:
            raise OfficeUploadError("uploaded presentation exceeds the shape count limit")
        slide_title, title_shape_index = _slide_title(slide)
        if builder.cursor:
            builder.append("\n\n")
        heading = f"# Slide {slide_number}"
        if slide_title:
            heading = f"{heading}: {_clean_inline_text(slide_title)}"
        heading_start, heading_end = builder.append(heading)
        _validate_markdown_budget(builder)
        elements.append(
            _source_element(
                kind="slide_heading",
                markdown_start=heading_start,
                markdown_end=heading_end,
                source_location={
                    "source_type": POWERPOINT_SOURCE_TYPE,
                    "slide_number": slide_number,
                    "slide_title": slide_title or None,
                    "shape_index": title_shape_index,
                },
            )
        )
        slide_has_body = False
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if getattr(shape, "has_table", False):
                table_cell_count += len(shape.table.rows) * len(shape.table.columns)
                if table_cell_count > MAX_POWERPOINT_TABLE_CELLS:
                    raise OfficeUploadError(
                        "uploaded presentation exceeds the table cell count limit"
                    )
                markdown = _shape_table_markdown(shape)
                if markdown is None:
                    continue
                builder.append("\n\n")
                start, end = builder.append(markdown)
                _validate_markdown_budget(builder)
                elements.append(
                    _source_element(
                        kind="slide_table",
                        markdown_start=start,
                        markdown_end=end,
                        source_location={
                            "source_type": POWERPOINT_SOURCE_TYPE,
                            "slide_number": slide_number,
                            "slide_title": slide_title or None,
                            "shape_index": shape_index,
                            "row_count": len(shape.table.rows),
                            "column_count": len(shape.table.columns),
                        },
                    )
                )
                slide_has_body = True
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            text = _shape_text(shape)
            if not text:
                continue
            if shape_index == title_shape_index and text == slide_title:
                continue
            builder.append("\n\n")
            start, end = builder.append(text)
            _validate_markdown_budget(builder)
            elements.append(
                _source_element(
                    kind="slide_text",
                    markdown_start=start,
                    markdown_end=end,
                    source_location={
                        "source_type": POWERPOINT_SOURCE_TYPE,
                        "slide_number": slide_number,
                        "slide_title": slide_title or None,
                        "shape_index": shape_index,
                    },
                )
            )
            slide_has_body = True
        if not slide_has_body and not slide_title:
            warnings.append(f"slide {slide_number} has no extractable text or tables")

    markdown_content = builder.content()
    if not markdown_content:
        raise OfficeUploadError("uploaded presentation has no extractable text or tables")
    return ParsedOfficeDocument(
        content=markdown_content,
        source_type=POWERPOINT_SOURCE_TYPE,
        source_content_type=POWERPOINT_CONTENT_TYPE,
        byte_size=len(content),
        sha256=hashlib.sha256(content).hexdigest(),
        parser_name=POWERPOINT_UPLOAD_PARSER_NAME,
        parser_version=_package_version("python-pptx"),
        elements=elements,
        warnings=warnings,
    )


def parse_uploaded_docx(
    *,
    filename: str | None,
    content_type: str | None,
    content: bytes,
    docling_config: DoclingExtractionConfig | None = None,
) -> ParsedOfficeDocument:
    """Parse a local `.docx` document into canonical Markdown and block locations."""
    safe_filename = _validate_office_filename(filename, expected_suffix=_WORD_DOCX_SUFFIX)
    _validate_content_type(
        suffix=_WORD_DOCX_SUFFIX,
        content_type=content_type,
        allowed_content_types=_WORD_DOCX_CONTENT_TYPES,
        expected_content_type=WORD_DOCX_CONTENT_TYPE,
    )
    _validate_office_content(content, filename=safe_filename)
    _validate_ooxml_archive(content, filename=safe_filename)
    markdown_content, warnings = _docx_markdown_with_docling(
        filename=safe_filename,
        content=content,
        config=docling_config or DoclingExtractionConfig(),
    )
    if not markdown_content:
        raise OfficeUploadError("uploaded Word document has no extractable text")
    return ParsedOfficeDocument(
        content=markdown_content,
        source_type=WORD_DOCUMENT_SOURCE_TYPE,
        source_content_type=WORD_DOCX_CONTENT_TYPE,
        byte_size=len(content),
        sha256=hashlib.sha256(content).hexdigest(),
        parser_name=WORD_DOCX_UPLOAD_PARSER_NAME,
        parser_version=_package_version("docling"),
        elements=_docx_markdown_elements(markdown_content),
        warnings=warnings,
    )


def _validate_office_filename(
    filename: str | None,
    *,
    expected_suffix: str | None = None,
) -> str:
    if filename is None or not filename.strip():
        raise OfficeUploadError("Office upload requires a filename")
    stripped = filename.strip()
    if "/" in stripped or "\\" in stripped:
        raise OfficeUploadError("Office filename must not contain path separators")
    suffix = _filename_suffix(stripped)
    if expected_suffix is not None and suffix != expected_suffix:
        raise OfficeUploadError(f"only {expected_suffix} uploads are supported")
    if expected_suffix is None and suffix not in _SUPPORTED_OFFICE_SUFFIXES:
        raise OfficeUploadError(
            f"only {_SUPPORTED_OFFICE_SUFFIX_LIST} Office uploads are supported"
        )
    return stripped


def _validate_content_type(
    *,
    suffix: str,
    content_type: str | None,
    allowed_content_types: frozenset[str],
    expected_content_type: str,
) -> None:
    normalized = _normalize_content_type(content_type)
    if normalized in allowed_content_types:
        return
    raise OfficeUploadError(f"{suffix} uploads must use {expected_content_type} compatible content")


def _validate_office_content(content: bytes, *, filename: str) -> None:
    if not content:
        raise OfficeUploadError("uploaded Office file is empty")
    if len(content) > MAX_OFFICE_UPLOAD_BYTES:
        raise OfficeUploadError("uploaded Office file exceeds the 5 MiB V1 limit")
    if not content.startswith(b"PK"):
        raise OfficeUploadError(f"{filename} is not a valid OOXML Office file")


def _validate_ooxml_archive(content: bytes, *, filename: str) -> None:
    try:
        with zipfile.ZipFile(BytesIO(content)) as archive:
            members = archive.infolist()
    except zipfile.BadZipFile as exc:
        raise OfficeUploadError(f"{filename} is not a valid OOXML Office file") from exc
    if len(members) > MAX_OFFICE_ARCHIVE_MEMBERS:
        raise OfficeUploadError("uploaded Office file exceeds the archive member count limit")
    total_uncompressed = 0
    for member in members:
        _validate_archive_member(member)
        total_uncompressed += max(member.file_size, 0)
        if total_uncompressed > MAX_OFFICE_ARCHIVE_UNCOMPRESSED_BYTES:
            raise OfficeUploadError("uploaded Office file exceeds the expanded size limit")
    compressed_size = max(len(content), 1)
    if total_uncompressed / compressed_size > MAX_OFFICE_ARCHIVE_COMPRESSION_RATIO:
        raise OfficeUploadError("uploaded Office file exceeds the compression ratio limit")


def _validate_archive_member(member: zipfile.ZipInfo) -> None:
    path = member.filename.replace("\\", "/")
    if (
        not path
        or path in {".", ".."}
        or path.startswith("/")
        or path.startswith("../")
        or "/../" in f"/{path}"
        or path.endswith("/..")
    ):
        raise OfficeUploadError("uploaded Office file contains an unsafe archive path")


def _docx_markdown_with_docling(
    *,
    filename: str,
    content: bytes,
    config: DoclingExtractionConfig,
) -> tuple[str, list[str]]:
    try:
        from docling.datamodel.accelerator_options import AcceleratorDevice, AcceleratorOptions
        from docling.datamodel.base_models import DocumentStream, InputFormat
        from docling.datamodel.document import ConversionStatus
        from docling.datamodel.pipeline_options import ConvertPipelineOptions
        from docling.document_converter import DocumentConverter, WordFormatOption
    except Exception as exc:  # noqa: BLE001 - local parser dependency should fail closed.
        raise OfficeUploadError("DOCX parsing is unavailable in this environment") from exc

    try:
        source = DocumentStream(name=filename, stream=BytesIO(content))
        pipeline_options = ConvertPipelineOptions()
        pipeline_options.accelerator_options = AcceleratorOptions(
            num_threads=config.threads,
            device=_docling_accelerator_device(config.accelerator, AcceleratorDevice),
        )
        pipeline_options.document_timeout = config.timeout_seconds
        converter = DocumentConverter(
            allowed_formats=[InputFormat.DOCX],
            format_options={InputFormat.DOCX: WordFormatOption(pipeline_options=pipeline_options)},
        )
        result = converter.convert(
            source,
            raises_on_error=False,
            max_file_size=MAX_OFFICE_UPLOAD_BYTES,
        )
    except Exception as exc:  # noqa: BLE001 - keep parser failures display-safe.
        raise OfficeUploadError(f"{filename} could not be parsed as a valid .docx file") from exc

    document = getattr(result, "document", None)
    status = getattr(result, "status", None)
    accepted_statuses = {ConversionStatus.SUCCESS, ConversionStatus.PARTIAL_SUCCESS}
    if status not in accepted_statuses or document is None:
        detail = _docling_error_summary(getattr(result, "errors", ()))
        suffix = f": {detail}" if detail else ""
        raise OfficeUploadError(f"{filename} could not be parsed as a valid .docx file{suffix}")

    markdown = _normalize_docx_markdown(document.export_to_markdown())
    if len(markdown) > MAX_OFFICE_MARKDOWN_CHARS:
        raise OfficeUploadError("uploaded Office file exceeds the extracted text limit")
    warnings = _docling_conversion_warnings(
        status_name=getattr(status, "name", str(status)),
        errors=getattr(result, "errors", ()),
        partial=status == ConversionStatus.PARTIAL_SUCCESS,
    )
    return markdown, warnings


def _normalize_docx_markdown(markdown: str) -> str:
    return markdown.replace("\r\n", "\n").replace("\r", "\n").strip()


def _docx_markdown_elements(markdown: str) -> list[dict[str, Any]]:
    lines = markdown.splitlines(keepends=True)
    line_offsets: list[tuple[int, str]] = []
    cursor = 0
    for line in lines:
        line_offsets.append((cursor, line))
        cursor += len(line)

    elements: list[dict[str, Any]] = []
    heading_stack: list[str] = []
    block_index = 0
    index = 0
    while index < len(line_offsets):
        line_start, line = line_offsets[index]
        stripped = line.strip()
        if not stripped:
            index += 1
            continue

        heading_match = _MARKDOWN_HEADING_PATTERN.match(stripped)
        if heading_match is not None:
            block_index += 1
            level = len(heading_match.group("marks"))
            title = _clean_inline_text(heading_match.group("title"))
            heading_stack = [*heading_stack[: level - 1], title]
            line_end = _line_markdown_end(line_start=line_start, line=line)
            elements.append(
                _word_source_element(
                    kind="word_heading",
                    markdown_start=line_start,
                    markdown_end=line_end,
                    block_index=block_index,
                    heading_path=heading_stack,
                    extra={"heading_level": level, "heading": title},
                )
            )
            index += 1
            continue

        if _is_markdown_table_line(stripped):
            block_index += 1
            table_start = line_start
            table_end, row_count, index = _consume_markdown_table(line_offsets, index)
            elements.append(
                _word_source_element(
                    kind="word_table",
                    markdown_start=table_start,
                    markdown_end=table_end,
                    block_index=block_index,
                    heading_path=heading_stack,
                    extra={"row_count": row_count},
                )
            )
            continue

        if _is_markdown_list_item(stripped):
            block_index += 1
            line_end = _line_markdown_end(line_start=line_start, line=line)
            elements.append(
                _word_source_element(
                    kind="word_list_item",
                    markdown_start=line_start,
                    markdown_end=line_end,
                    block_index=block_index,
                    heading_path=heading_stack,
                    extra={},
                )
            )
            index += 1
            continue

        block_index += 1
        paragraph_start = line_start
        paragraph_end, index = _consume_markdown_paragraph(line_offsets, index)
        elements.append(
            _word_source_element(
                kind="word_paragraph",
                markdown_start=paragraph_start,
                markdown_end=paragraph_end,
                block_index=block_index,
                heading_path=heading_stack,
                extra={},
            )
        )

    return elements


def _consume_markdown_table(
    line_offsets: Sequence[tuple[int, str]],
    index: int,
) -> tuple[int, int, int]:
    table_end = line_offsets[index][0]
    row_count = 0
    while index < len(line_offsets):
        line_start, line = line_offsets[index]
        stripped = line.strip()
        if not _is_markdown_table_line(stripped):
            break
        table_end = _line_markdown_end(line_start=line_start, line=line)
        if not _is_markdown_table_separator(stripped):
            row_count += 1
        index += 1
    return table_end, row_count, index


def _consume_markdown_paragraph(
    line_offsets: Sequence[tuple[int, str]],
    index: int,
) -> tuple[int, int]:
    paragraph_end = line_offsets[index][0]
    while index < len(line_offsets):
        line_start, line = line_offsets[index]
        stripped = line.strip()
        if (
            not stripped
            or _MARKDOWN_HEADING_PATTERN.match(stripped)
            or _is_markdown_table_line(stripped)
            or _is_markdown_list_item(stripped)
        ):
            break
        paragraph_end = _line_markdown_end(line_start=line_start, line=line)
        index += 1
    return paragraph_end, index


def _word_source_element(
    *,
    kind: str,
    markdown_start: int,
    markdown_end: int,
    block_index: int,
    heading_path: list[str],
    extra: dict[str, Any],
) -> dict[str, Any]:
    source_location = {
        "source_type": WORD_DOCUMENT_SOURCE_TYPE,
        "block_index": block_index,
        "heading_path": " > ".join(heading_path) if heading_path else None,
        **extra,
    }
    return _source_element(
        kind=kind,
        markdown_start=markdown_start,
        markdown_end=markdown_end,
        source_location=source_location,
    )


def _line_markdown_end(*, line_start: int, line: str) -> int:
    return line_start + len(line.rstrip("\n"))


def _is_markdown_table_line(stripped_line: str) -> bool:
    return stripped_line.startswith("|") and stripped_line.endswith("|")


def _is_markdown_table_separator(stripped_line: str) -> bool:
    cells = [cell.strip() for cell in stripped_line.strip("|").split("|")]
    return bool(cells) and all(cell and set(cell) <= {"-", ":"} for cell in cells)


def _is_markdown_list_item(stripped_line: str) -> bool:
    return stripped_line.startswith(("- ", "* ", "+ ")) or bool(
        _MARKDOWN_ORDERED_LIST_PATTERN.match(stripped_line)
    )


def _docling_conversion_warnings(
    *,
    status_name: str,
    errors: Sequence[Any],
    partial: bool,
) -> list[str]:
    warnings: list[str] = []
    if partial:
        warnings.append(f"docling conversion status: {status_name}")
    warnings.extend(_safe_docling_message(error) for error in errors[:5])
    return [warning for warning in warnings if warning]


def _docling_error_summary(errors: Sequence[Any]) -> str:
    return "; ".join(_safe_docling_message(error) for error in errors[:3] if error)[:200]


def _safe_docling_message(error: Any) -> str:
    message = getattr(error, "message", None) or str(error)
    return _clean_inline_text(message)[:200]


def _worksheet_markdown(worksheet: Any) -> _WorksheetMarkdown | None:
    grid = _worksheet_grid(worksheet)
    if grid is None:
        return None

    width = max((len(row) for _, row in grid.rows), default=0)
    header_row_count = _worksheet_header_row_count(
        grid.rows,
        merged_cell_count=grid.merged_cell_count,
    )
    header_rows = [row for _, row in grid.rows[:header_row_count]]
    body_rows = grid.rows[header_row_count:]
    header = _flatten_header_rows(header_rows, width=width)
    header = _apply_spreadsheet_header_hints(header, [row for _, row in body_rows])
    table = _markdown_table(header, [_pad_row(row, width=width) for _, row in body_rows])
    first_row = grid.rows[0][0]
    last_row = grid.rows[-1][0]
    first_column = get_column_letter(grid.first_column_index)
    last_column = get_column_letter(grid.last_column_index)
    source_location = {
        "source_type": EXCEL_SOURCE_TYPE,
        "sheet_name": worksheet.title,
        "cell_range": f"{first_column}{first_row}:{last_column}{last_row}",
        "row_start": first_row,
        "row_end": last_row,
        "column_start": first_column,
        "column_end": last_column,
    }

    builder = _MarkdownBuilder()
    builder.append(f"## Sheet: {_clean_inline_text(worksheet.title)}\n\n{table}")
    elements: list[dict[str, Any]] = []
    warnings: list[str] = []
    if _should_emit_row_records(
        merged_cell_count=grid.merged_cell_count,
        header_row_count=header_row_count,
    ):
        builder.append("\n\n### Row records")
        emitted_records = 0
        for row_index, row in body_rows:
            if emitted_records >= MAX_EXCEL_ROW_RECORDS_PER_SHEET:
                warnings.append(
                    f"worksheet {worksheet.title!r} row records were truncated at "
                    f"{MAX_EXCEL_ROW_RECORDS_PER_SHEET} rows"
                )
                break
            record = _worksheet_row_record(row_index=row_index, header=header, row=row)
            if record is None:
                continue
            builder.append("\n\n")
            start, end = builder.append(record)
            emitted_records += 1
            elements.append(
                _source_element(
                    kind="worksheet_row_record",
                    markdown_start=start,
                    markdown_end=end,
                    source_location={
                        "source_type": EXCEL_SOURCE_TYPE,
                        "sheet_name": worksheet.title,
                        "cell_range": f"{first_column}{row_index}:{last_column}{row_index}",
                        "row_start": row_index,
                        "row_end": row_index,
                        "column_start": first_column,
                        "column_end": last_column,
                    },
                )
            )

    markdown = builder.content()
    elements.insert(
        0,
        _source_element(
            kind="worksheet_table",
            markdown_start=0,
            markdown_end=len(markdown),
            source_location=source_location,
        ),
    )
    return _WorksheetMarkdown(markdown=markdown, elements=elements, warnings=warnings)


def _worksheet_grid(worksheet: Any) -> _WorksheetGrid | None:
    rows: list[tuple[int, list[str]]] = []
    first_column_index: int | None = None
    last_column_index = 0
    non_empty_cell_count = 0
    merged_values = _merged_cell_values(worksheet)
    merged_cell_count = len(getattr(getattr(worksheet, "merged_cells", None), "ranges", ()))
    raw_rows: list[tuple[int, list[str]]] = []
    for row_index, cells in enumerate(worksheet.iter_rows(), start=1):
        if row_index > MAX_EXCEL_ROWS_PER_SHEET:
            raise OfficeUploadError("uploaded spreadsheet exceeds the row count limit")
        normalized = [
            _cell_to_markdown(merged_values.get((row_index, cell.column), cell.value))
            for cell in cells
        ]
        non_empty_cell_count += sum(1 for cell in cells if _cell_to_markdown(cell.value))
        if non_empty_cell_count > MAX_EXCEL_NON_EMPTY_CELLS:
            raise OfficeUploadError("uploaded spreadsheet exceeds the cell count limit")
        trimmed, row_first_column, row_last_column = _trim_empty_cells(normalized)
        if not trimmed:
            continue
        raw_rows.append((row_index, normalized))
        if row_first_column is not None:
            first_column_index = (
                row_first_column
                if first_column_index is None
                else min(first_column_index, row_first_column)
            )
            last_column_index = max(last_column_index, row_last_column)
    if not raw_rows or first_column_index is None:
        return None
    rows = [
        (row_index, values[first_column_index - 1 : last_column_index])
        for row_index, values in raw_rows
    ]
    return _WorksheetGrid(
        rows=rows,
        first_column_index=first_column_index,
        last_column_index=last_column_index,
        merged_cell_count=merged_cell_count,
    )


def _merged_cell_values(worksheet: Any) -> dict[tuple[int, int], Any]:
    ranges = getattr(getattr(worksheet, "merged_cells", None), "ranges", ())
    merged_values: dict[tuple[int, int], Any] = {}
    filled_cells = 0
    for merged_range in ranges:
        value = worksheet.cell(merged_range.min_row, merged_range.min_col).value
        if value is None:
            continue
        row_end = min(merged_range.max_row, MAX_EXCEL_ROWS_PER_SHEET)
        for row_index in range(merged_range.min_row, row_end + 1):
            for column_index in range(merged_range.min_col, merged_range.max_col + 1):
                filled_cells += 1
                if filled_cells > MAX_EXCEL_NON_EMPTY_CELLS:
                    raise OfficeUploadError("uploaded spreadsheet exceeds the cell count limit")
                merged_values[(row_index, column_index)] = value
    return merged_values


def _worksheet_header_row_count(
    rows: list[tuple[int, list[str]]],
    *,
    merged_cell_count: int,
) -> int:
    if len(rows) < 2:
        return 1
    first_row = rows[0][1]
    might_have_layered_header = merged_cell_count > 0 or any(not cell for cell in first_row)
    if not might_have_layered_header:
        return 1
    header_row_count = 1
    for _, row in rows[1 : min(len(rows), 5)]:
        if _looks_like_spreadsheet_data_row(row):
            break
        header_row_count += 1
    return header_row_count


def _looks_like_spreadsheet_data_row(row: list[str]) -> bool:
    values = [value for value in row if value]
    if not values:
        return False
    numeric_count = sum(1 for value in values if _looks_like_numeric_cell(value))
    month_count = sum(1 for value in values if _MONTH_VALUE_PATTERN.fullmatch(value))
    return numeric_count >= max(1, len(values) // 4) or bool(month_count and numeric_count)


def _looks_like_numeric_cell(value: str) -> bool:
    normalized = _clean_inline_text(value).replace(",", "")
    return bool(_NUMERIC_VALUE_PATTERN.fullmatch(normalized))


def _flatten_header_rows(header_rows: list[list[str]], *, width: int) -> list[str]:
    flattened: list[str] = []
    padded_rows = [_pad_row(row, width=width) for row in header_rows]
    for column_index in range(width):
        parts: list[str] = []
        seen_parts: set[str] = set()
        for row in padded_rows:
            part = _clean_inline_text(row[column_index])
            if not part or part.casefold() in seen_parts:
                continue
            seen_parts.add(part.casefold())
            parts.append(part)
        flattened.append(" ".join(parts))
    return _normalized_header(flattened, width=width)


def _apply_spreadsheet_header_hints(
    header: list[str],
    body_rows: list[list[str]],
) -> list[str]:
    if not header or not body_rows:
        return header
    first_column_values = [row[0] for row in body_rows if row and row[0]]
    month_count = sum(1 for value in first_column_values if _MONTH_VALUE_PATTERN.fullmatch(value))
    if month_count >= 2 and month_count * 2 >= len(first_column_values):
        hinted = [*header]
        hinted[0] = "월"
        return _normalized_header(hinted, width=len(hinted))
    return header


def _should_emit_row_records(
    *,
    merged_cell_count: int,
    header_row_count: int,
) -> bool:
    return merged_cell_count > 0 or header_row_count > 1


def _worksheet_row_record(
    *,
    row_index: int,
    header: list[str],
    row: list[str],
) -> str | None:
    pairs: list[str] = []
    for key, value in zip(header, _pad_row(row, width=len(header)), strict=True):
        if not value:
            continue
        pairs.append(f"{_clean_inline_text(key)}={_clean_inline_text(value)}")
    if not pairs:
        return None
    return f"- Excel row {row_index}: " + "; ".join(pairs)


def _shape_table_markdown(shape: Any) -> str | None:
    table = shape.table
    rows = [[_cell_to_markdown(cell.text) for cell in row.cells] for row in table.rows]
    rows = [_pad_row(row, width=max((len(row) for row in rows), default=0)) for row in rows]
    rows = [row for row in rows if any(cell for cell in row)]
    if not rows:
        return None
    width = max(len(row) for row in rows)
    header = _normalized_header(rows[0], width=width)
    body_rows = [_pad_row(row, width=width) for row in rows[1:]]
    return _markdown_table(header, body_rows)


def _slide_title(slide: Any) -> tuple[str, int | None]:
    for shape_index, shape in enumerate(slide.shapes, start=1):
        if not getattr(shape, "has_text_frame", False):
            continue
        if getattr(shape, "is_placeholder", False) and getattr(
            shape.placeholder_format, "type", None
        ):
            text = _shape_text(shape)
            if text:
                return text, shape_index
    for shape_index, shape in enumerate(slide.shapes, start=1):
        if not getattr(shape, "has_text_frame", False):
            continue
        text = _shape_text(shape)
        if text:
            return text, shape_index
    return "", None


def _shape_text(shape: Any) -> str:
    text = getattr(shape, "text", "")
    return _normalize_text(text)


def _normalized_header(row: list[str], *, width: int) -> list[str]:
    padded = _pad_row(row, width=width)
    headers: list[str] = []
    seen: set[str] = set()
    for index, value in enumerate(padded, start=1):
        header = value or f"Column {get_column_letter(index)}"
        unique_header = header
        suffix = 2
        while unique_header.casefold() in seen:
            unique_header = f"{header} {suffix}"
            suffix += 1
        seen.add(unique_header.casefold())
        headers.append(unique_header)
    return headers


def _markdown_table(header: list[str], body_rows: list[list[str]]) -> str:
    rows = [header, ["---" for _ in header], *body_rows]
    return "\n".join(
        "| " + " | ".join(_escape_markdown_cell(cell) for cell in row) + " |" for row in rows
    )


def _trim_empty_cells(values: list[str]) -> tuple[list[str], int | None, int]:
    first = next((index for index, value in enumerate(values) if value), None)
    if first is None:
        return [], None, 0
    last = len(values) - 1
    while last >= first and not values[last]:
        last -= 1
    return values[first : last + 1], first + 1, last + 1


def _pad_row(row: list[str], *, width: int) -> list[str]:
    return [*row, *("" for _ in range(max(0, width - len(row))))]


def _cell_to_markdown(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, bool):
        return "TRUE" if value else "FALSE"
    text = str(value)
    return _normalize_text(text)


def _normalize_text(text: str) -> str:
    return "\n".join(
        line.strip()
        for line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n")
        if line.strip()
    )


def _escape_markdown_cell(value: str) -> str:
    return value.replace("\\", "\\\\").replace("|", "\\|").replace("\n", " ")


def _clean_inline_text(value: str) -> str:
    return _normalize_text(value).replace("\n", " ")


def _source_element(
    *,
    kind: str,
    markdown_start: int,
    markdown_end: int,
    source_location: dict[str, Any],
) -> dict[str, Any]:
    return {
        "kind": kind,
        "markdown_start": markdown_start,
        "markdown_end": markdown_end,
        "source_location": source_location,
    }


def _validate_markdown_budget(builder: _MarkdownBuilder) -> None:
    if builder.cursor > MAX_OFFICE_MARKDOWN_CHARS:
        raise OfficeUploadError("uploaded Office file exceeds the extracted text limit")


def _filename_suffix(filename: str) -> str:
    if "." not in filename:
        return ""
    return f".{filename.rsplit('.', maxsplit=1)[-1].casefold()}"


def _normalize_content_type(content_type: str | None) -> str:
    return (content_type or "").split(";", maxsplit=1)[0].strip().casefold()


def _package_version(package_name: str) -> str | None:
    try:
        return version(package_name)
    except PackageNotFoundError:
        return None
