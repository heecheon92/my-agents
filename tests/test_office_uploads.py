"""Unit tests for deterministic Office upload parsers."""

from __future__ import annotations

from io import BytesIO
from typing import Any, cast

import pytest
from docx import Document as WordDocument
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

import my_agents.knowledge.office_uploads as office_uploads_module
from my_agents.knowledge.extraction import _chunk_text
from my_agents.knowledge.office_uploads import (
    EXCEL_CONTENT_TYPE,
    EXCEL_UPLOAD_PARSER_NAME,
    POWERPOINT_CONTENT_TYPE,
    POWERPOINT_UPLOAD_PARSER_NAME,
    WORD_DOCUMENT_SOURCE_TYPE,
    WORD_DOCX_CONTENT_TYPE,
    WORD_DOCX_UPLOAD_PARSER_NAME,
    OfficeUploadError,
    parse_uploaded_office_document,
)
from my_agents.knowledge.pdf_uploads import DoclingExtractionConfig


def _xlsx_bytes() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Pipeline"
    sheet.append(["Metric", "Value"])
    sheet.append(["Latency", 120])
    sheet.append(["Owner", "Agent Platform"])
    empty = workbook.create_sheet("Empty Sheet")
    empty["A1"] = None
    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _merged_header_xlsx_bytes() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "원주시 PM10"
    sheet.append(
        [
            "시,도명",
            "도시명",
            "측정소명",
            "유효자료\n획득율(%)",
            "유효\n측정일수",
            "유효\n측정시간",
            "월평균\n(㎍/㎥)",
            "24시간치",
            None,
            None,
            None,
            None,
        ]
    )
    sheet.append(
        [
            None,
            None,
            None,
            None,
            None,
            None,
            None,
            "최저\n(㎍/㎥)",
            "최고\n(㎍/㎥)",
            "최고일시\n(년월일시)",
            "기준초과\n(회)",
            "초과율\n(%)",
        ]
    )
    sheet.append(["12월", "원주시", "중앙동", 98.79, 31, 735, 43, 14, 88, 20201211, 0, 0])
    sheet.append([None, None, "반곡동", 99.46, 31, 740, 43, 18, 85, 20201211, 0, 0])
    sheet.append([None, None, "문막읍", 99.06, 31, 737, 39, 18, 73, 20201211, 0, 0])
    sheet.append([None, None, "도시평균", 99.1, 93, 2212, 42, 14, 88, 20201211, 0, 0])
    for cell_range in ("A1:A2", "B1:B2", "C1:C2", "D1:D2", "E1:E2", "F1:F2", "G1:G2"):
        sheet.merge_cells(cell_range)
    sheet.merge_cells("H1:L1")
    sheet.merge_cells("A3:A6")
    sheet.merge_cells("B3:B6")
    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Upload Plan"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.4), Inches(6), Inches(0.8))
    textbox.text = "Support PPTX parsing"
    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2.4), Inches(5), Inches(1.2))
    table = table_shape.table
    table.cell(0, 0).text = "Phase"
    table.cell(0, 1).text = "Owner"
    table.cell(1, 0).text = "Parser"
    table.cell(1, 1).text = "Backend"
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _docx_bytes() -> bytes:
    document = WordDocument()
    document.add_heading("DOCX Upload Plan", level=1)
    document.add_paragraph("WordAlpha supports DOCX ingestion.")
    document.add_paragraph("GET /documents/upload")
    document.add_paragraph("First bullet", style="List Bullet")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "DocxParser"
    table.cell(1, 1).text = "Ready"
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def test_parse_xlsx_to_markdown_with_source_location_offsets() -> None:
    content = _xlsx_bytes()

    parsed = parse_uploaded_office_document(
        filename="metrics.xlsx",
        content_type=EXCEL_CONTENT_TYPE,
        content=content,
    )

    assert parsed.source_type == "spreadsheet"
    assert parsed.source_content_type == EXCEL_CONTENT_TYPE
    assert parsed.parser_name == EXCEL_UPLOAD_PARSER_NAME
    assert parsed.sha256
    assert parsed.content == (
        "## Sheet: Pipeline\n\n"
        "| Metric | Value |\n"
        "| --- | --- |\n"
        "| Latency | 120 |\n"
        "| Owner | Agent Platform |"
    )
    assert parsed.warnings == ["worksheet 'Empty Sheet' has no extractable cells"]

    [element] = parsed.elements
    assert element["kind"] == "worksheet_table"
    assert parsed.content[element["markdown_start"] : element["markdown_end"]] == parsed.content
    assert element["source_location"] == {
        "source_type": "spreadsheet",
        "sheet_name": "Pipeline",
        "cell_range": "A1:B3",
        "row_start": 1,
        "row_end": 3,
        "column_start": "A",
        "column_end": "B",
    }


def test_parse_xlsx_expands_merged_headers_into_row_records() -> None:
    content = _merged_header_xlsx_bytes()

    parsed = parse_uploaded_office_document(
        filename="wonju-pm10.xlsx",
        content_type=EXCEL_CONTENT_TYPE,
        content=content,
    )

    assert "Column I" not in parsed.content
    assert (
        "| 월 | 도시명 | 측정소명 | 유효자료 획득율(%) | 유효 측정일수 | "
        "유효 측정시간 | 월평균 (㎍/㎥) | 24시간치 최저 (㎍/㎥) | "
        "24시간치 최고 (㎍/㎥) | 24시간치 최고일시 (년월일시) | "
        "24시간치 기준초과 (회) | 24시간치 초과율 (%) |"
    ) in parsed.content
    assert (
        "| 12월 | 원주시 | 도시평균 | 99.1 | 93 | 2212 | 42 | 14 | 88 | 20201211 | 0 | 0 |"
    ) in parsed.content
    city_average_record = (
        "- Excel row 6: 월=12월; 도시명=원주시; 측정소명=도시평균; "
        "유효자료 획득율(%)=99.1; 유효 측정일수=93; 유효 측정시간=2212; "
        "월평균 (㎍/㎥)=42; 24시간치 최저 (㎍/㎥)=14; "
        "24시간치 최고 (㎍/㎥)=88; 24시간치 최고일시 (년월일시)=20201211; "
        "24시간치 기준초과 (회)=0; 24시간치 초과율 (%)=0"
    )
    assert city_average_record in parsed.content

    row_record_element = next(
        element
        for element in parsed.elements
        if element["kind"] == "worksheet_row_record"
        and element["source_location"]["row_start"] == 6
    )
    assert (
        parsed.content[row_record_element["markdown_start"] : row_record_element["markdown_end"]]
        == city_average_record
    )
    assert row_record_element["source_location"] == {
        "source_type": "spreadsheet",
        "sheet_name": "원주시 PM10",
        "cell_range": "A6:L6",
        "row_start": 6,
        "row_end": 6,
        "column_start": "A",
        "column_end": "L",
    }

    chunk_texts = [chunk for chunk, _, _ in _chunk_text(parsed.content)]
    matching_chunks = [chunk for chunk in chunk_texts if "Excel row 6" in chunk]
    assert matching_chunks == [city_average_record]


def test_parse_pptx_to_markdown_with_slide_and_shape_locations() -> None:
    content = _pptx_bytes()

    parsed = parse_uploaded_office_document(
        filename="deck.pptx",
        content_type=POWERPOINT_CONTENT_TYPE,
        content=content,
    )

    assert parsed.source_type == "presentation"
    assert parsed.source_content_type == POWERPOINT_CONTENT_TYPE
    assert parsed.parser_name == POWERPOINT_UPLOAD_PARSER_NAME
    assert "# Slide 1: Upload Plan" in parsed.content
    assert "Support PPTX parsing" in parsed.content
    assert "| Phase | Owner |" in parsed.content

    locations = [element["source_location"] for element in parsed.elements]
    assert {location["slide_number"] for location in locations} == {1}
    assert any(element["kind"] == "slide_heading" for element in parsed.elements)
    text_element = next(element for element in parsed.elements if element["kind"] == "slide_text")
    assert parsed.content[text_element["markdown_start"] : text_element["markdown_end"]] == (
        "Support PPTX parsing"
    )
    assert text_element["source_location"]["slide_title"] == "Upload Plan"
    assert text_element["source_location"]["shape_index"] is not None


def test_parse_docx_to_markdown_with_word_block_locations() -> None:
    content = _docx_bytes()

    parsed = parse_uploaded_office_document(
        filename="word-plan.docx",
        content_type=WORD_DOCX_CONTENT_TYPE,
        content=content,
    )

    assert parsed.source_type == WORD_DOCUMENT_SOURCE_TYPE
    assert parsed.source_content_type == WORD_DOCX_CONTENT_TYPE
    assert parsed.parser_name == WORD_DOCX_UPLOAD_PARSER_NAME
    assert parsed.parser_provider == "local"
    assert parsed.parser_version
    assert parsed.sha256
    assert parsed.content.startswith("## DOCX Upload Plan")
    assert "WordAlpha supports DOCX ingestion." in parsed.content
    assert "GET /documents/upload" in parsed.content
    assert "DocxParser" in parsed.content
    assert "Ready" in parsed.content

    kinds = {element["kind"] for element in parsed.elements}
    assert "word_heading" in kinds
    assert "word_paragraph" in kinds
    assert "word_table" in kinds

    heading = next(element for element in parsed.elements if element["kind"] == "word_heading")
    assert parsed.content[heading["markdown_start"] : heading["markdown_end"]] == (
        "## DOCX Upload Plan"
    )
    assert heading["source_location"] == {
        "source_type": WORD_DOCUMENT_SOURCE_TYPE,
        "block_index": 1,
        "heading_path": "DOCX Upload Plan",
        "heading_level": 2,
        "heading": "DOCX Upload Plan",
    }

    endpoint = next(
        element
        for element in parsed.elements
        if parsed.content[element["markdown_start"] : element["markdown_end"]]
        == "GET /documents/upload"
    )
    assert endpoint["kind"] == "word_paragraph"
    assert endpoint["source_location"] == {
        "source_type": WORD_DOCUMENT_SOURCE_TYPE,
        "block_index": 3,
        "heading_path": "DOCX Upload Plan",
    }

    table = next(element for element in parsed.elements if element["kind"] == "word_table")
    assert "| Metric" in parsed.content[table["markdown_start"] : table["markdown_end"]]
    assert table["source_location"] == {
        "source_type": WORD_DOCUMENT_SOURCE_TYPE,
        "block_index": 5,
        "heading_path": "DOCX Upload Plan",
        "row_count": 2,
    }


def test_docx_parser_uses_docling_extraction_config(monkeypatch) -> None:  # noqa: ANN001
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.document import ConversionStatus

    captured: dict[str, object] = {}

    class _FakeDocument:
        def export_to_markdown(self) -> str:
            return "## Configured DOCX\n\nDoclingTimeoutAlpha"

    class _FakeResult:
        document = _FakeDocument()
        status = ConversionStatus.SUCCESS
        errors: tuple[object, ...] = ()

    class _FakeConverter:
        def __init__(self, *args: object, **kwargs: object) -> None:
            captured["format_options"] = kwargs["format_options"]

        def convert(self, *args: object, **kwargs: object) -> _FakeResult:
            return _FakeResult()

    monkeypatch.setattr("docling.document_converter.DocumentConverter", _FakeConverter)

    parsed = parse_uploaded_office_document(
        filename="word-plan.docx",
        content_type=WORD_DOCX_CONTENT_TYPE,
        content=_docx_bytes(),
        docling_config=DoclingExtractionConfig(
            accelerator="cpu",
            ocr_enabled=False,
            timeout_seconds=12.5,
            threads=2,
        ),
    )

    format_options = cast(dict[object, Any], captured["format_options"])
    word_options = format_options[InputFormat.DOCX]
    pipeline_options = word_options.pipeline_options
    assert pipeline_options.document_timeout == 12.5
    assert pipeline_options.accelerator_options.num_threads == 2
    device = pipeline_options.accelerator_options.device
    assert getattr(device, "value", device) == "cpu"
    assert "DoclingTimeoutAlpha" in parsed.content


def test_docx_parser_rejects_extracted_markdown_over_budget(monkeypatch) -> None:  # noqa: ANN001
    from docling.datamodel.document import ConversionStatus

    class _FakeDocument:
        def export_to_markdown(self) -> str:
            return "A" * 33

    class _FakeResult:
        document = _FakeDocument()
        status = ConversionStatus.SUCCESS
        errors: tuple[object, ...] = ()

    class _FakeConverter:
        def __init__(self, *args: object, **kwargs: object) -> None:
            pass

        def convert(self, *args: object, **kwargs: object) -> _FakeResult:
            return _FakeResult()

    monkeypatch.setattr(office_uploads_module, "MAX_OFFICE_MARKDOWN_CHARS", 32)
    monkeypatch.setattr("docling.document_converter.DocumentConverter", _FakeConverter)

    with pytest.raises(OfficeUploadError, match="extracted text limit"):
        parse_uploaded_office_document(
            filename="word-plan.docx",
            content_type=WORD_DOCX_CONTENT_TYPE,
            content=_docx_bytes(),
        )


def test_office_parser_rejects_unsupported_or_corrupted_files_safely() -> None:
    with pytest.raises(OfficeUploadError, match=r"only \.xlsx, \.pptx, and \.docx"):
        parse_uploaded_office_document(
            filename="legacy.xls",
            content_type="application/vnd.ms-excel",
            content=b"not-office",
        )

    with pytest.raises(OfficeUploadError, match="not a valid OOXML Office file"):
        parse_uploaded_office_document(
            filename="metrics.xlsx",
            content_type=EXCEL_CONTENT_TYPE,
            content=b"not-office",
        )

    with pytest.raises(OfficeUploadError, match="not a valid OOXML Office file"):
        parse_uploaded_office_document(
            filename="word-plan.docx",
            content_type=WORD_DOCX_CONTENT_TYPE,
            content=b"not-office",
        )

    with pytest.raises(OfficeUploadError, match=r"\.docx uploads must use"):
        parse_uploaded_office_document(
            filename="word-plan.docx",
            content_type="text/plain",
            content=_docx_bytes(),
        )

    with pytest.raises(OfficeUploadError, match=r"only \.xlsx, \.pptx, and \.docx"):
        parse_uploaded_office_document(
            filename="legacy.doc",
            content_type="application/msword",
            content=b"not-office",
        )


def test_office_parser_rejects_archive_expansion_before_parser_load(monkeypatch) -> None:  # noqa: ANN001
    monkeypatch.setattr(office_uploads_module, "MAX_OFFICE_ARCHIVE_UNCOMPRESSED_BYTES", 1)

    with pytest.raises(OfficeUploadError, match="expanded size limit"):
        parse_uploaded_office_document(
            filename="metrics.xlsx",
            content_type=EXCEL_CONTENT_TYPE,
            content=_xlsx_bytes(),
        )


def test_office_parser_rejects_spreadsheets_over_cell_budget(monkeypatch) -> None:  # noqa: ANN001
    monkeypatch.setattr(office_uploads_module, "MAX_EXCEL_NON_EMPTY_CELLS", 1)

    with pytest.raises(OfficeUploadError, match="cell count limit"):
        parse_uploaded_office_document(
            filename="metrics.xlsx",
            content_type=EXCEL_CONTENT_TYPE,
            content=_xlsx_bytes(),
        )


def test_office_parser_rejects_presentations_over_shape_budget(monkeypatch) -> None:  # noqa: ANN001
    monkeypatch.setattr(office_uploads_module, "MAX_POWERPOINT_SHAPES", 1)

    with pytest.raises(OfficeUploadError, match="shape count limit"):
        parse_uploaded_office_document(
            filename="deck.pptx",
            content_type=POWERPOINT_CONTENT_TYPE,
            content=_pptx_bytes(),
        )
