"""Knowledge-base ingestion and deterministic extraction tests."""

from __future__ import annotations

import json
import logging
import zlib
from io import BytesIO
from pathlib import Path
from time import monotonic, sleep

import pytest
from fastapi.testclient import TestClient
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
from sqlalchemy import select

import my_agents.knowledge.extraction as extraction_module
import my_agents.knowledge.pdf_uploads as pdf_uploads_module
from my_agents.api.conversations.retrieval_context import retrieved_context_for_graph
from my_agents.knowledge.extraction import (
    _chunk_pdf_text,
    _chunk_text,
    _deterministic_embedding,
    _extract_entity_names,
)
from my_agents.knowledge.ingestion_worker import process_pending_extraction_runs_once
from my_agents.knowledge.models import (
    CitationModel,
    DocumentChunkModel,
    DocumentModel,
    DocumentParseArtifactModel,
    DocumentPermissionModel,
    EntityMentionModel,
    EntityRelationshipModel,
    ExtractionRunModel,
    KnowledgeBaseModel,
    StructuredKnowledgeEntityModel,
)
from my_agents.knowledge.pdf_uploads import PdfUploadError, parse_uploaded_pdf
from my_agents.knowledge.retrieval import RetrievalService
from my_agents.persistence.database import get_database_session

from .conftest import load_app, verify_latest_auth_email


class FakeEmbeddingProvider:
    provider = "fake"
    model = "fake-embedding-model"
    dimensions = 3

    def embed_documents(self, texts: list[str]) -> list[list[float]]:
        return [[float(index), float(len(text)), 1.0] for index, text in enumerate(texts)]

    def embed_query(self, text: str) -> list[float]:
        return [1.0, 0.0, 0.0]


class FailingEmbeddingProvider:
    provider = "failing"
    model = "failing-embedding-model"
    dimensions = 3

    def embed_documents(self, texts: list[str]) -> list[list[float]]:  # noqa: ARG002
        raise RuntimeError("fixture embedding failure")

    def embed_query(self, text: str) -> list[float]:  # noqa: ARG002
        return [1.0, 0.0, 0.0]


class SlowEmbeddingProvider(FakeEmbeddingProvider):
    provider = "slow-fake"
    model = "slow-fake-embedding-model"

    def embed_documents(self, texts: list[str]) -> list[list[float]]:
        sleep(0.05)
        return super().embed_documents(texts)


def _client(monkeypatch) -> TestClient:  # noqa: ANN001 - pytest monkeypatch fixture
    monkeypatch.setenv("MY_AGENTS_RESPONSE_MODE", "deterministic")
    monkeypatch.setenv("MY_AGENTS_SESSION_COOKIE_SECURE", "false")
    return TestClient(load_app())


def _file_client(monkeypatch, tmp_path) -> TestClient:  # noqa: ANN001 - pytest fixtures
    monkeypatch.setenv("MY_AGENTS_RESPONSE_MODE", "deterministic")
    monkeypatch.setenv("MY_AGENTS_SESSION_COOKIE_SECURE", "false")
    monkeypatch.setenv("MY_AGENTS_DATABASE_URL", f"sqlite+pysqlite:///{tmp_path / 'app.db'}")
    monkeypatch.setenv("MY_AGENTS_AUTO_CREATE_TABLES", "true")
    return TestClient(load_app())


def _signup_login(client: TestClient, email: str) -> str:
    password = "correct horse battery staple"
    signup = client.post("/auth/signup", json={"email": email, "password": password})
    assert signup.status_code == 201
    verify_latest_auth_email(client, email)
    login = client.post("/auth/login", json={"email": email, "password": password})
    assert login.status_code == 200
    return signup.json()["user"]["id"]


def _create_knowledge_base(client: TestClient, name: str = "Test KB") -> str:
    response = client.post("/knowledge-bases", json={"name": name, "scope": "personal"})
    assert response.status_code == 201
    return response.json()["id"]


def _create_personal_knowledge_base(client: TestClient, name: str = "Test KB") -> str:
    return _create_knowledge_base(client, name)


def _with_knowledge_base(client: TestClient, payload: dict) -> dict:
    if "knowledge_base_id" in payload:
        return payload
    return {**payload, "knowledge_base_id": _create_knowledge_base(client)}


def _create_document(client: TestClient, *, json: dict):  # noqa: ANN201
    return client.post("/documents", json=_with_knowledge_base(client, json))


def _upload_document(client: TestClient, *, data: dict, files: dict):  # noqa: ANN201
    upload_data = dict(data)
    upload_data.setdefault("knowledge_base_id", _create_knowledge_base(client))
    return client.post("/documents/upload", data=upload_data, files=files)


def _text_pdf(*pages: str) -> bytes:
    objects = ["1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj"]
    kids = []
    for index, page in enumerate(pages, start=1):
        page_obj = 2 + (index * 2) - 1
        stream_obj = page_obj + 1
        kids.append(f"{page_obj} 0 R")
        escaped = page.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
        stream = f"BT /F1 12 Tf 72 720 Td ({escaped}) Tj ET"
        objects.append(
            f"{page_obj} 0 obj << /Type /Page /Parent 2 0 R /Contents {stream_obj} 0 R >> endobj"
        )
        objects.append(
            f"{stream_obj} 0 obj << /Length {len(stream)} >> stream\n{stream}\nendstream endobj"
        )
    objects.insert(
        1, f"2 0 obj << /Type /Pages /Kids [{' '.join(kids)}] /Count {len(kids)} >> endobj"
    )
    return ("%PDF-1.4\n" + "\n".join(objects) + "\n%%EOF\n").encode()


def _compressed_text_pdf(*pages: str) -> bytes:
    objects = ["1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj"]
    kids = []
    for index, page in enumerate(pages, start=1):
        page_obj = 2 + (index * 2) - 1
        stream_obj = page_obj + 1
        kids.append(f"{page_obj} 0 R")
        escaped = page.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
        stream = f"BT /F1 12 Tf 72 720 Td ({escaped}) Tj ET".encode()
        compressed = zlib.compress(stream)
        objects.append(
            f"{page_obj} 0 obj << /Type /Page /Parent 2 0 R /Contents {stream_obj} 0 R >> endobj"
        )
        objects.append(
            b"".join(
                [
                    (
                        f"{stream_obj} 0 obj "
                        f"<< /Length {len(compressed)} /Filter /FlateDecode >> stream\n"
                    ).encode(),
                    compressed,
                    b"\nendstream endobj",
                ]
            ).decode("latin-1")
        )
    objects.insert(
        1, f"2 0 obj << /Type /Pages /Kids [{' '.join(kids)}] /Count {len(kids)} >> endobj"
    )
    return ("%PDF-1.4\n" + "\n".join(objects) + "\n%%EOF\n").encode("latin-1")


def _binary_noise_pdf() -> bytes:
    noise = b"$\xa6\xedO\x7f /\x89\xe1\x88e<m\xde(\x8d\x97Xz\xf7x\xb5q;)\xff\xfe"
    objects = [
        b"1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj",
        b"2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj",
        b"3 0 obj << /Type /Page /Parent 2 0 R /Contents 4 0 R >> endobj",
        b"4 0 obj << /Length "
        + str(len(noise)).encode()
        + b" >> stream\n"
        + noise
        + b"\nendstream endobj",
    ]
    return b"%PDF-1.4\n" + b"\n".join(objects) + b"\n%%EOF\n"


def _raw_stream_pdf(stream: str) -> bytes:
    objects = [
        "1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj",
        "2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj",
        "3 0 obj << /Type /Page /Parent 2 0 R /Contents 4 0 R >> endobj",
        f"4 0 obj << /Length {len(stream)} >> stream\n{stream}\nendstream endobj",
    ]
    return ("%PDF-1.4\n" + "\n".join(objects) + "\n%%EOF\n").encode()


def _xlsx_bytes(*, sheet_name: str = "Pipeline") -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = sheet_name
    sheet.append(["Metric", "Value"])
    sheet.append(["Latency", 120])
    sheet.append(["Owner", "Agent Platform"])
    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _provenance_xlsx_bytes() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Provenance"
    sheet.append(["Signal", "Endpoint"])
    sheet.append(["OfficeProvenanceAlpha", "GET /office-metrics"])
    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Office Slide Plan"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.4), Inches(6), Inches(0.8))
    textbox.text = "OfficeSlideAlpha depends on PowerPoint provenance."
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _database_rows(statement):  # noqa: ANN001 - SQLAlchemy statement type is verbose
    session_generator = get_database_session()
    db = next(session_generator)
    try:
        return db.scalars(statement).all()
    finally:
        session_generator.close()


def _wait_for_run(
    client: TestClient,
    document_id: str,
    run_id: str,
    *,
    terminal_status: str = "completed",
) -> dict:
    deadline = monotonic() + 10
    payload: dict | None = None
    while monotonic() < deadline:
        response = client.get(f"/documents/{document_id}/extraction-runs/{run_id}")
        assert response.status_code == 200
        payload = response.json()
        if payload["status"] == terminal_status:
            return payload
        sleep(0.02)
    pytest.fail(f"run {run_id} did not reach {terminal_status}; last payload={payload}")


def test_personal_knowledge_base_document_ingestion_creates_extraction_artifacts(
    monkeypatch,
) -> None:  # noqa: ANN001
    client = _client(monkeypatch)
    _signup_login(client, "kb-owner@example.com")

    kb = client.post("/knowledge-bases", json={"name": "Personal KB", "scope": "personal"})
    assert kb.status_code == 201
    kb_id = kb.json()["id"]

    document = _create_document(
        client,
        json={
            "title": "Agent Notes",
            "content": "OpenAI builds agents with LangGraph.\n\nLangGraph helps Heecheon Park.",
            "knowledge_base_id": kb_id,
        },
    )
    assert document.status_code == 201
    assert document.json()["knowledge_base_id"] == kb_id
    assert document.json()["source_type"] == "text"
    assert document.json()["source_filename"] is None

    ingest = client.post(f"/documents/{document.json()['id']}/ingest")

    assert ingest.status_code == 200
    payload = ingest.json()
    assert payload["status"] == "completed"
    assert payload["chunk_count"] == 2
    assert payload["entity_count"] >= 3
    assert payload["relationship_count"] >= 1

    runs = client.get(f"/documents/{document.json()['id']}/extraction-runs")
    assert runs.status_code == 200
    assert runs.json()[0]["id"] == payload["id"]
    chunks = _database_rows(
        select(DocumentChunkModel).where(DocumentChunkModel.document_id == document.json()["id"])
    )
    assert {chunk.source_page for chunk in chunks} == {None}


def test_long_text_chunking_uses_document_sized_target_and_overlap() -> None:
    text = "A" * 1800

    chunks = _chunk_text(text)

    assert [(start, end) for _, start, end in chunks] == [(0, 1500), (1300, 1800)]
    assert [len(content) for content, *_ in chunks] == [1500, 500]


def test_reingesting_document_replaces_chunks_without_duplicate_ordinals(
    monkeypatch,
) -> None:  # noqa: ANN001
    client = _client(monkeypatch)
    _signup_login(client, "kb-reingest-owner@example.com")
    kb_id = _create_knowledge_base(client, "Reingest KB")
    document = _create_document(
        client,
        json={
            "title": "Reingest Notes",
            "content": "First reingest paragraph.\n\nSecond reingest paragraph.",
            "knowledge_base_id": kb_id,
        },
    )
    assert document.status_code == 201
    document_id = document.json()["id"]

    first = client.post(f"/documents/{document_id}/ingest")
    second = client.post(f"/documents/{document_id}/ingest")

    assert first.status_code == 200
    assert second.status_code == 200
    assert first.json()["chunk_count"] == 2
    assert second.json()["chunk_count"] == 2
    chunks = _database_rows(
        select(DocumentChunkModel)
        .where(DocumentChunkModel.document_id == document_id)
        .order_by(DocumentChunkModel.ordinal)
    )

    assert [chunk.ordinal for chunk in chunks] == [0, 1]
    assert [chunk.content for chunk in chunks] == [
        "First reingest paragraph.",
        "Second reingest paragraph.",
    ]


def test_legacy_document_create_without_knowledge_base_is_rejected(monkeypatch) -> None:  # noqa: ANN001
    client = _client(monkeypatch)
    _signup_login(client, "missing-kb-create@example.com")

    response = client.post(
        "/documents",
        json={"title": "Legacy Missing KB", "content": "must not persist without a KB"},
    )

    assert response.status_code == 422
    assert _database_rows(select(DocumentModel)) == []


def test_legacy_document_upload_without_knowledge_base_is_rejected(monkeypatch) -> None:  # noqa: ANN001
    client = _client(monkeypatch)
    _signup_login(client, "missing-kb-upload@example.com")

    response = client.post(
        "/documents/upload",
        data={"title": "Legacy Upload Missing KB"},
        files={"file": ("missing-kb.txt", b"missing knowledge base", "text/plain")},
    )

    assert response.status_code == 422
    assert _database_rows(select(DocumentModel)) == []


@pytest.mark.parametrize("endpoint", ["/documents", "/documents/upload"])
def test_document_write_with_nonexistent_knowledge_base_returns_404(
    monkeypatch,
    endpoint: str,
) -> None:  # noqa: ANN001
    client = _client(monkeypatch)
    endpoint_name = endpoint.rsplit("/", maxsplit=1)[-1] or "create"
    _signup_login(client, f"missing-kb-{endpoint_name}@example.com")
    missing_kb_id = "00000000-0000-0000-0000-000000000000"

    if endpoint == "/documents":
        response = client.post(
            endpoint,
            json={
                "title": "Missing KB",
                "content": "no document should be stored",
                "knowledge_base_id": missing_kb_id,
            },
        )
    else:
        response = client.post(
            endpoint,
            data={"title": "Missing KB Upload", "knowledge_base_id": missing_kb_id},
            files={"file": ("missing-kb.txt", b"missing knowledge base", "text/plain")},
        )

    assert response.status_code == 404
    assert response.json()["detail"] == "knowledge base not found"
    assert _database_rows(select(DocumentModel)) == []


@pytest.mark.parametrize("endpoint", ["/documents", "/documents/upload"])
def test_document_write_with_unauthorized_knowledge_base_returns_404(
    monkeypatch,
    endpoint: str,
) -> None:  # noqa: ANN001
    owner = _client(monkeypatch)
    outsider = _client(monkeypatch)
    endpoint_name = endpoint.rsplit("/", maxsplit=1)[-1] or "create"
    _signup_login(owner, f"kb-auth-owner-{endpoint_name}@example.com")
    _signup_login(
        outsider,
        f"kb-auth-outsider-{endpoint_name}@example.com",
    )
    kb_id = _create_personal_knowledge_base(owner, "Owner-only KB")

    if endpoint == "/documents":
        response = outsider.post(
            endpoint,
            json={
                "title": "Unauthorized KB",
                "content": "no document should be stored",
                "knowledge_base_id": kb_id,
            },
        )
    else:
        response = outsider.post(
            endpoint,
            data={"title": "Unauthorized KB Upload", "knowledge_base_id": kb_id},
            files={"file": ("unauthorized-kb.txt", b"unauthorized", "text/plain")},
        )

    assert response.status_code == 404
    assert response.json()["detail"] == "knowledge base not found"
    assert _database_rows(select(DocumentModel)) == []


def test_document_write_paths_do_not_create_null_knowledge_base_ids(monkeypatch) -> None:  # noqa: ANN001
    client = _client(monkeypatch)
    _signup_login(client, "no-null-kb@example.com")
    text_kb_id = _create_personal_knowledge_base(client, "Text No Null KB")
    upload_kb_id = _create_personal_knowledge_base(client, "Upload No Null KB")

    created = client.post(
        "/documents",
        json={
            "title": "No Null Text",
            "content": "text create must carry a KB",
            "knowledge_base_id": text_kb_id,
        },
    )
    uploaded = client.post(
        "/documents/upload",
        data={"title": "No Null Upload", "knowledge_base_id": upload_kb_id},
        files={"file": ("no-null.txt", b"upload create must carry a KB", "text/plain")},
    )

    assert created.status_code == 201
    assert uploaded.status_code == 201
    documents = _database_rows(select(DocumentModel))
    assert {document.knowledge_base_id for document in documents} == {text_kb_id, upload_kb_id}
    null_kb_documents = _database_rows(
        select(DocumentModel).where(DocumentModel.knowledge_base_id.is_(None))
    )
    assert null_kb_documents == []
    assert _database_rows(select(KnowledgeBaseModel)) != []


def test_ingestion_uses_configured_embedding_provider(monkeypatch) -> None:  # noqa: ANN001
    monkeypatch.setattr(
        extraction_module,
        "get_embedding_provider",
        lambda: FakeEmbeddingProvider(),
    )
    client = _client(monkeypatch)
    _signup_login(client, "provider-ingest@example.com")
    kb_id = _create_personal_knowledge_base(client, "Provider KB")

    document = _create_document(
        client,
        json={
            "title": "Provider Embeddings",
            "content": "First provider chunk.\n\nSecond provider chunk.",
            "knowledge_base_id": kb_id,
        },
    )
    assert document.status_code == 201

    ingest = client.post(f"/documents/{document.json()['id']}/ingest")

    assert ingest.status_code == 200
    chunks = _database_rows(
        select(DocumentChunkModel)
        .where(DocumentChunkModel.document_id == document.json()["id"])
        .order_by(DocumentChunkModel.ordinal)
    )
    assert [chunk.embedding_json for chunk in chunks] == [
        "[0.0, 21.0, 1.0]",
        "[1.0, 22.0, 1.0]",
    ]


def test_async_ingest_returns_queued_run_and_polling_completes(monkeypatch) -> None:  # noqa: ANN001
    client = _client(monkeypatch)
    _signup_login(client, "async-ingest-owner@example.com")
    kb_id = _create_personal_knowledge_base(client, "Async KB")

    document = _create_document(
        client,
        json={
            "title": "Async Notes",
            "content": "Async ingestion mentions LangGraph.\n\nSecond chunk mentions FastAPI.",
            "knowledge_base_id": kb_id,
        },
    )
    assert document.status_code == 201
    document_id = document.json()["id"]

    queued = client.post(f"/documents/{document_id}/ingest/async")

    assert queued.status_code == 202
    queued_payload = queued.json()
    assert queued_payload["document_id"] == document_id
    assert queued_payload["status"] == "pending"
    assert queued_payload["stage"] == "queued"
    assert queued_payload["progress_percent"] == 0
    assert queued_payload["error"] is None

    completed = _wait_for_run(client, document_id, queued_payload["id"])
    assert completed["status"] == "completed"
    assert completed["stage"] == "completed"
    assert completed["progress_percent"] == 100
    assert completed["chunk_count"] >= 1
    assert completed["entity_count"] >= 2
    assert completed["relationship_count"] >= 1
    assert completed["started_at"] is not None
    assert completed["completed_at"] is not None

    runs = client.get(f"/documents/{document_id}/extraction-runs")
    assert runs.status_code == 200
    assert runs.json()[0]["id"] == queued_payload["id"]


def test_async_ingest_can_be_processed_by_external_worker(monkeypatch) -> None:  # noqa: ANN001
    monkeypatch.setenv("MY_AGENTS_INGESTION_EXECUTION_MODE", "external_worker")
    client = _client(monkeypatch)
    _signup_login(client, "external-worker-owner@example.com")
    kb_id = _create_personal_knowledge_base(client, "External Worker KB")

    document = _create_document(
        client,
        json={
            "title": "Worker Notes",
            "content": "External worker ingestion mentions LangGraph and FastAPI.",
            "knowledge_base_id": kb_id,
        },
    )
    assert document.status_code == 201
    document_id = document.json()["id"]

    queued = client.post(f"/documents/{document_id}/ingest/async")

    assert queued.status_code == 202
    queued_payload = queued.json()
    assert queued_payload["status"] == "pending"
    assert queued_payload["stage"] == "queued"
    before_worker = client.get(f"/documents/{document_id}/extraction-runs/{queued_payload['id']}")
    assert before_worker.status_code == 200
    assert before_worker.json()["status"] == "pending"

    results = process_pending_extraction_runs_once(
        database_url="sqlite+pysqlite:///:memory:",
        max_runs=1,
    )

    assert len(results) == 1
    assert results[0].run_id == queued_payload["id"]
    assert results[0].status == "completed"
    completed = client.get(f"/documents/{document_id}/extraction-runs/{queued_payload['id']}")
    assert completed.status_code == 200
    assert completed.json()["status"] == "completed"
    assert completed.json()["chunk_count"] >= 1


def test_async_ingest_persists_failed_status_with_safe_error(monkeypatch) -> None:  # noqa: ANN001
    monkeypatch.setattr(
        extraction_module,
        "get_embedding_provider",
        lambda: FailingEmbeddingProvider(),
    )
    client = _client(monkeypatch)
    _signup_login(client, "async-ingest-failure@example.com")
    kb_id = _create_personal_knowledge_base(client, "Async Failure KB")

    document = _create_document(
        client,
        json={
            "title": "Failed Async",
            "content": "This run should fail embedding.",
            "knowledge_base_id": kb_id,
        },
    )
    assert document.status_code == 201
    document_id = document.json()["id"]

    queued = client.post(f"/documents/{document_id}/ingest/async")

    assert queued.status_code == 202
    failed = _wait_for_run(
        client,
        document_id,
        queued.json()["id"],
        terminal_status="failed",
    )
    assert failed["status"] == "failed"
    assert failed["stage"] == "failed"
    assert "RuntimeError" in failed["error"]
    assert "fixture embedding failure" in failed["error"]


def test_async_ingest_start_and_poll_respect_document_permissions(monkeypatch) -> None:  # noqa: ANN001
    owner = _client(monkeypatch)
    reader = _client(monkeypatch)
    outsider = _client(monkeypatch)
    _signup_login(owner, "async-permission-owner@example.com")
    reader_id = _signup_login(reader, "async-permission-reader@example.com")
    _signup_login(outsider, "async-permission-outsider@example.com")
    kb_id = _create_personal_knowledge_base(owner, "Async Permission KB")

    document = _create_document(
        owner,
        json={
            "title": "Permission Async",
            "content": "Reader can poll but not ingest.",
            "knowledge_base_id": kb_id,
        },
    )
    assert document.status_code == 201
    document_id = document.json()["id"]
    grant = owner.patch(
        f"/documents/{document_id}/permissions",
        json={"user_id": reader_id, "can_read": True, "can_ingest": False},
    )
    assert grant.status_code == 200

    denied_start = reader.post(f"/documents/{document_id}/ingest/async")
    assert denied_start.status_code == 403

    queued = owner.post(f"/documents/{document_id}/ingest/async")
    assert queued.status_code == 202
    run_id = queued.json()["id"]
    _wait_for_run(owner, document_id, run_id)

    reader_poll = reader.get(f"/documents/{document_id}/extraction-runs/{run_id}")
    assert reader_poll.status_code == 200
    outsider_poll = outsider.get(f"/documents/{document_id}/extraction-runs/{run_id}")
    assert outsider_poll.status_code == 404


def test_parallel_async_ingest_shared_entities_complete(monkeypatch, tmp_path) -> None:  # noqa: ANN001
    monkeypatch.setattr(
        extraction_module,
        "get_embedding_provider",
        lambda: SlowEmbeddingProvider(),
    )
    client = _file_client(monkeypatch, tmp_path)
    _signup_login(client, "parallel-async-owner@example.com")
    kb_id = _create_personal_knowledge_base(client, "Parallel Async KB")

    document_ids = []
    for index, content in enumerate(
        [
            "Shared Alpha uses FastAPI. LangGraph helps Shared Alpha.",
            "FastAPI helps Shared Alpha. Shared Alpha studies LangGraph.",
            "LangGraph and FastAPI support Shared Alpha in public demos.",
        ],
        start=1,
    ):
        document = _create_document(
            client,
            json={
                "title": f"Parallel Async {index}",
                "content": content,
                "knowledge_base_id": kb_id,
            },
        )
        assert document.status_code == 201
        document_ids.append(document.json()["id"])

    queued_runs = [
        client.post(f"/documents/{document_id}/ingest/async") for document_id in document_ids
    ]

    assert [response.status_code for response in queued_runs] == [202, 202, 202]
    completed_runs = [
        _wait_for_run(client, document_id, response.json()["id"])
        for document_id, response in zip(document_ids, queued_runs, strict=True)
    ]

    assert [run["status"] for run in completed_runs] == ["completed", "completed", "completed"]
    assert all(run["entity_count"] >= 2 for run in completed_runs)


def test_pdf_parser_tolerates_invalid_octal_like_literal_escape() -> None:
    parsed = parse_uploaded_pdf(
        filename="invalid-escape.pdf",
        content_type="application/pdf",
        content=_raw_stream_pdf(r"BT /F1 12 Tf 72 720 Td (Invalid \9 escape) Tj ET"),
    )

    assert parsed.content == "Invalid 9 escape"
    assert parsed.page_count == 1


def test_pdf_parser_decodes_flate_streams() -> None:
    parsed = parse_uploaded_pdf(
        filename="compressed.pdf",
        content_type="application/pdf",
        content=_compressed_text_pdf("Compressed resume mentions FastAPI and LangGraph."),
    )

    assert parsed.content == "Compressed resume mentions FastAPI and LangGraph."
    assert parsed.page_count == 1


def test_pdf_parser_removes_postgres_unsafe_nul_bytes() -> None:
    parsed = parse_uploaded_pdf(
        filename="nul-byte.pdf",
        content_type="application/pdf",
        content=_raw_stream_pdf(r"BT /F1 12 Tf 72 720 Td (Alpha \000 Beta) Tj ET"),
    )

    assert parsed.content == "Alpha Beta"
    assert "\x00" not in parsed.content
    assert parsed.page_count == 1


def test_pdf_parser_uses_docling_after_pymupdf_quality_failure(monkeypatch) -> None:  # noqa: ANN001
    monkeypatch.setattr(pdf_uploads_module, "_extract_page_texts_with_pymupdf", lambda _: [])
    monkeypatch.setattr(pdf_uploads_module, "_extract_page_texts_with_pypdf", lambda _: [])
    monkeypatch.setattr(
        pdf_uploads_module,
        "_extract_page_texts_with_docling",
        lambda filename, _, config: [  # noqa: ARG005
            f"Docling extracted {filename} as Markdown."
        ],
    )

    parsed = parse_uploaded_pdf(
        filename="docling.pdf",
        content_type="application/pdf",
        content=_text_pdf("PyMuPDF would normally parse this."),
    )

    assert parsed.content == "Docling extracted docling.pdf as Markdown."
    assert parsed.parser_name == "docling_markdown_v1"


def test_docling_extractor_forces_cpu_accelerator(monkeypatch) -> None:  # noqa: ANN001
    captured: dict[str, object] = {}

    class FakeDocument:
        def num_pages(self) -> int:
            return 1

        def export_to_markdown(self, *, page_no: int | None = None) -> str:  # noqa: ARG002
            return "Docling CPU output mentions FastAPI."

    class FakeConverter:
        def __init__(self, *, format_options: dict) -> None:
            pdf_options = next(iter(format_options.values())).pipeline_options
            captured["device"] = pdf_options.accelerator_options.device
            captured["do_ocr"] = pdf_options.do_ocr
            captured["timeout"] = pdf_options.document_timeout

        def convert(self, *args, **kwargs):  # noqa: ANN002, ANN003
            return type("Result", (), {"document": FakeDocument()})()

    monkeypatch.setattr(
        pdf_uploads_module,
        "_extract_page_texts_with_pymupdf",
        lambda _: [],
    )
    monkeypatch.setattr(pdf_uploads_module, "_extract_page_texts_with_pypdf", lambda _: [])
    monkeypatch.setattr(
        "docling.document_converter.DocumentConverter",
        FakeConverter,
    )

    parsed = parse_uploaded_pdf(
        filename="docling-cpu.pdf",
        content_type="application/pdf",
        content=_text_pdf("force docling path"),
    )

    assert parsed.parser_name == "docling_markdown_v1"
    assert parsed.content == "Docling CPU output mentions FastAPI."
    assert str(captured["device"]).endswith("CPU")
    assert captured["do_ocr"] is False
    assert captured["timeout"] == 30.0


def test_pdf_parser_uses_tesseract_after_docling_quality_failure(monkeypatch) -> None:  # noqa: ANN001
    monkeypatch.setattr(pdf_uploads_module, "_extract_page_texts_with_pymupdf", lambda _: [])
    monkeypatch.setattr(pdf_uploads_module, "_extract_page_texts_with_pypdf", lambda _: [])
    monkeypatch.setattr(
        pdf_uploads_module,
        "_extract_page_texts_with_docling",
        lambda filename, content, docling_config: ["<!-- image -->"],  # noqa: ARG005
    )
    monkeypatch.setattr(
        pdf_uploads_module,
        "_extract_page_texts_with_tesseract",
        lambda content, config: [f"Tesseract OCR with {config.languages} found Korean text."],
    )

    parsed = parse_uploaded_pdf(
        filename="ocr.pdf",
        content_type="application/pdf",
        content=_text_pdf("image-heavy fixture"),
    )

    assert parsed.content == "Tesseract OCR with kor+eng found Korean text."
    assert parsed.parser_name == "tesseract_ocr_v1"


def test_tesseract_extractor_returns_empty_when_disabled() -> None:
    pages = pdf_uploads_module._extract_page_texts_with_tesseract(
        _text_pdf("disabled OCR fixture"),
        pdf_uploads_module.TesseractOcrConfig(enabled=False),
    )

    assert pages == []


def test_tesseract_extractor_skips_pdfs_over_max_page_cap(monkeypatch) -> None:  # noqa: ANN001
    monkeypatch.setattr(pdf_uploads_module.shutil, "which", lambda _: "/usr/bin/tesseract")
    monkeypatch.setattr(
        pdf_uploads_module,
        "_run_tesseract_page_ocr",
        lambda *_, **__: pytest.fail("large PDFs must not spawn per-page OCR"),
    )

    pages = pdf_uploads_module._extract_page_texts_with_tesseract(
        _text_pdf("first page", "second page"),
        pdf_uploads_module.TesseractOcrConfig(enabled=True, max_pages=1),
    )

    assert pages == []


def test_pdf_parser_rejects_docling_image_placeholders_without_chunks(monkeypatch) -> None:  # noqa: ANN001
    monkeypatch.setattr(pdf_uploads_module, "_extract_page_texts_with_pymupdf", lambda _: [])
    monkeypatch.setattr(
        pdf_uploads_module,
        "_extract_page_texts_with_docling",
        lambda filename, content, config: [  # noqa: ARG005
            "<!-- image -->\n<!-- image -->",
            "- ▪\n- ▪\n- ▪",
            "<!-- image -->",
        ],
    )
    monkeypatch.setattr(pdf_uploads_module, "_extract_page_texts_with_pypdf", lambda _: [])
    monkeypatch.setattr(pdf_uploads_module, "_extract_page_texts_with_tesseract", lambda *_: [])
    monkeypatch.setattr(pdf_uploads_module, "_legacy_extract_page_texts", lambda _: [])

    with pytest.raises(PdfUploadError, match="does not contain extractable text"):
        parse_uploaded_pdf(
            filename="image-placeholder.pdf",
            content_type="application/pdf",
            content=_text_pdf("image placeholder fixture"),
        )


def test_pdf_parser_allows_table_of_contents_dot_leaders() -> None:
    repeated_dots = "Vision" + "." * 80 + "3\nMission" + "." * 80 + "4"

    parsed = parse_uploaded_pdf(
        filename="toc-dot-leaders.pdf",
        content_type="application/pdf",
        content=_text_pdf(repeated_dots),
    )

    assert parsed.parser_name == "pymupdf_text_v1"
    assert "Vision" in parsed.content
    assert "Mission" in parsed.content


def test_pdf_parser_rejects_repeated_non_punctuation_artifacts(monkeypatch) -> None:  # noqa: ANN001
    artifact = "A" * 80
    monkeypatch.setattr(
        pdf_uploads_module, "_extract_page_texts_with_pymupdf", lambda _: [artifact]
    )
    monkeypatch.setattr(pdf_uploads_module, "_extract_page_texts_with_pypdf", lambda _: [artifact])
    monkeypatch.setattr(
        pdf_uploads_module,
        "_extract_page_texts_with_docling",
        lambda filename, content, config: [artifact],  # noqa: ARG005
    )
    monkeypatch.setattr(pdf_uploads_module, "_extract_page_texts_with_tesseract", lambda *_: [])
    monkeypatch.setattr(pdf_uploads_module, "_legacy_extract_page_texts", lambda _: [artifact])

    with pytest.raises(PdfUploadError, match="repeated character artifacts"):
        parse_uploaded_pdf(
            filename="repeat-artifact.pdf",
            content_type="application/pdf",
            content=_text_pdf("source fixture text"),
        )


def test_nara_2022_pdf_regression_accepts_native_text_when_available() -> None:
    sample = Path(
        "/Users/heecheonpark/Downloads/KG-RAG-datasets-main/us-fed-agency-reports/data/v1/docs/NARA2022.pdf"
    )
    if not sample.exists():
        pytest.skip("local NARA2022 sample PDF is not available")

    parsed = parse_uploaded_pdf(
        filename=sample.name,
        content_type="application/pdf",
        content=sample.read_bytes(),
    )

    assert parsed.parser_name == "pymupdf_text_v1"
    assert "NARA 2022" in parsed.content
    assert "Strategic Plan" in parsed.content


def test_pdf_parser_rejects_binary_literal_noise() -> None:
    with pytest.raises(PdfUploadError, match="does not contain extractable text"):
        parse_uploaded_pdf(
            filename="noise.pdf",
            content_type="application/pdf",
            content=_binary_noise_pdf(),
        )


def test_nested_pdf_upload_rejects_garbled_locale_artifact_without_500(
    monkeypatch,  # noqa: ANN001
    caplog: pytest.LogCaptureFixture,
) -> None:
    caplog.set_level(logging.INFO)
    client = _client(monkeypatch)
    _signup_login(client, "pdf-garble@example.com")
    kb_id = _create_personal_knowledge_base(client, "PDF Garble KB")
    locale_artifact = " ".join(["ko-KR"] * 30)
    garbled_pdf = _raw_stream_pdf(
        rf"BT /F1 12 Tf 72 720 Td ({locale_artifact} \000 Adobe UCS) Tj ET"
    )

    response = client.post(
        f"/knowledge-bases/{kb_id}/documents/upload",
        data={"title": "Garbled PDF"},
        files={"file": ("garbled.pdf", garbled_pdf, "application/pdf")},
    )

    assert response.status_code == 400
    assert "does not contain extractable text" in response.json()["detail"]
    assert any(
        "document_upload.received" in record.message and "garbled.pdf" in record.message
        for record in caplog.records
    )
    assert any(
        "pdf_upload.parser.rejected" in record.message and "garbled.pdf" in record.message
        for record in caplog.records
    )
    assert any(
        "document_upload.rejected" in record.message and "garbled.pdf" in record.message
        for record in caplog.records
    )


def test_pdf_upload_persists_metadata_and_ingests_page_provenance(monkeypatch) -> None:  # noqa: ANN001
    client = _client(monkeypatch)
    _signup_login(client, "pdf-owner@example.com")
    kb_id = _create_personal_knowledge_base(client, "PDF KB")

    document = _upload_document(
        client,
        data={"title": "Demo PDF", "knowledge_base_id": kb_id},
        files={
            "file": (
                "demo.pdf",
                _text_pdf(
                    "OpenAI Agents page one mentions LangGraph.",
                    "Heecheon Park page two cites FastAPI.",
                ),
                "application/pdf",
            )
        },
    )

    assert document.status_code == 201
    payload = document.json()
    assert payload["title"] == "Demo PDF"
    assert payload["source_type"] == "pdf"
    assert payload["source_filename"] == "demo.pdf"
    assert payload["source_content_type"] == "application/pdf"
    assert payload["source_byte_size"] > 0
    assert len(payload["source_sha256"]) == 64
    assert payload["source_page_count"] == 2
    assert payload["parser_name"] == "pymupdf_text_v1"

    persisted = _database_rows(select(DocumentModel).where(DocumentModel.id == payload["id"]))
    assert persisted[0].content.count("\f") == 1

    ingest = client.post(f"/documents/{payload['id']}/ingest")
    assert ingest.status_code == 200
    assert ingest.json()["status"] == "completed"
    assert ingest.json()["chunk_count"] == 2

    chunks = _database_rows(
        select(DocumentChunkModel)
        .where(DocumentChunkModel.document_id == payload["id"])
        .order_by(DocumentChunkModel.ordinal)
    )
    assert [chunk.source_page for chunk in chunks] == [1, 2]
    assert [chunk.source_location_json for chunk in chunks] == [None, None]
    assert "LangGraph" in chunks[0].content
    assert "FastAPI" in chunks[1].content
    assert len(_deterministic_embedding(chunks[0].content)) == 32


@pytest.mark.parametrize(
    ("filename", "content_type", "source_type", "source_content_type", "parser_name"),
    [
        ("demo-notes.md", "text/markdown", "markdown", "text/markdown", "utf8_markdown_v1"),
        ("demo-notes.txt", "text/plain", "text", "text/plain", "utf8_text_v1"),
    ],
)
def test_text_upload_persists_metadata_and_ingests_for_retrieval(
    monkeypatch,
    filename: str,
    content_type: str,
    source_type: str,
    source_content_type: str,
    parser_name: str,
) -> None:  # noqa: ANN001
    client = _client(monkeypatch)
    _signup_login(client, f"{source_type}-upload-owner@example.com")
    kb_id = _create_personal_knowledge_base(client, f"{source_type.title()} Upload KB")

    phrase = f"TextUpload {source_type} source says Heecheon Park uses LangGraph retrieval"
    document = _upload_document(
        client,
        data={"title": f"{source_type.title()} Upload", "knowledge_base_id": kb_id},
        files={
            "file": (
                filename,
                f"# Product notes\r\n\r\n{phrase} with FastAPI citations.\n".encode(),
                content_type,
            )
        },
    )

    assert document.status_code == 201
    payload = document.json()
    assert payload["source_type"] == source_type
    assert payload["source_filename"] == filename
    assert payload["source_content_type"] == source_content_type
    assert payload["source_byte_size"] > 0
    assert len(payload["source_sha256"]) == 64
    assert payload["source_page_count"] is None
    assert payload["parser_name"] == parser_name

    persisted = _database_rows(select(DocumentModel).where(DocumentModel.id == payload["id"]))
    assert "\r" not in persisted[0].content
    assert phrase in persisted[0].content

    ingest = client.post(f"/documents/{payload['id']}/ingest")
    assert ingest.status_code == 200
    assert ingest.json()["status"] == "completed"
    assert ingest.json()["chunk_count"] >= 1

    conversation = client.post("/conversations", json={"title": f"{source_type} RAG"})
    assert conversation.status_code == 201
    run = client.post(
        f"/conversations/{conversation.json()['id']}/runs",
        json={"message": f"What does my uploaded {source_type} file say about TextUpload?"},
    )

    assert run.status_code == 200
    run_payload = run.json()
    assert run_payload["citations"]
    assert run_payload["citations"][0]["document_id"] == payload["id"]
    assert run_payload["citations"][0]["source_filename"] == filename
    assert phrase in run_payload["citations"][0]["snippet"]
    assert not run_payload["reply"].startswith("Based on authorized document context:")


def test_office_upload_persists_parse_artifact_and_delete_cleans_it_up(monkeypatch) -> None:  # noqa: ANN001
    client = _client(monkeypatch)
    _signup_login(client, "office-upload-owner@example.com")
    kb_id = _create_personal_knowledge_base(client, "Office Upload KB")

    document = _upload_document(
        client,
        data={"title": "Metrics Workbook", "knowledge_base_id": kb_id},
        files={
            "file": (
                "metrics.xlsx",
                _xlsx_bytes(),
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
        },
    )

    assert document.status_code == 201
    payload = document.json()
    assert payload["source_type"] == "spreadsheet"
    assert payload["source_filename"] == "metrics.xlsx"
    assert payload["source_content_type"] == (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )
    assert payload["parser_name"] == "openpyxl_markdown_v1"
    assert payload["source_page_count"] is None

    persisted = _database_rows(select(DocumentModel).where(DocumentModel.id == payload["id"]))
    assert persisted[0].content.startswith("## Sheet: Pipeline")
    assert "| Latency | 120 |" in persisted[0].content

    [artifact] = _database_rows(
        select(DocumentParseArtifactModel).where(
            DocumentParseArtifactModel.document_id == payload["id"]
        )
    )
    assert artifact.source_sha256 == payload["source_sha256"]
    assert artifact.source_filename == "metrics.xlsx"
    assert artifact.source_type == "spreadsheet"
    assert artifact.parser_provider == "local"
    assert artifact.parser_name == "openpyxl_markdown_v1"
    assert artifact.markdown_content == persisted[0].content
    assert json.loads(artifact.elements_json)[0]["source_location"]["cell_range"] == "A1:B3"
    assert json.loads(artifact.warnings_json) == []

    deleted = client.delete(f"/documents/{payload['id']}")

    assert deleted.status_code == 204
    assert (
        _database_rows(
            select(DocumentParseArtifactModel).where(
                DocumentParseArtifactModel.document_id == payload["id"]
            )
        )
        == []
    )


def test_office_ingestion_exposes_source_location_on_chunks_debug_and_citations(
    monkeypatch,
) -> None:  # noqa: ANN001
    client = _client(monkeypatch)
    user_id = _signup_login(client, "office-provenance-owner@example.com")
    kb_id = _create_personal_knowledge_base(client, "Office Provenance KB")
    document = _upload_document(
        client,
        data={"title": "Provenance Workbook", "knowledge_base_id": kb_id},
        files={
            "file": (
                "provenance.xlsx",
                _provenance_xlsx_bytes(),
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
        },
    )
    assert document.status_code == 201
    document_id = document.json()["id"]

    ingest = client.post(f"/documents/{document_id}/ingest")

    assert ingest.status_code == 200
    chunks = _database_rows(
        select(DocumentChunkModel)
        .where(DocumentChunkModel.document_id == document_id)
        .order_by(DocumentChunkModel.ordinal)
    )
    chunk = next(item for item in chunks if "OfficeProvenanceAlpha" in item.content)
    source_location = json.loads(chunk.source_location_json)
    assert source_location == {
        "cell_range": "A1:B2",
        "column_end": "B",
        "column_start": "A",
        "row_end": 2,
        "row_start": 1,
        "sheet_name": "Provenance",
        "source_type": "spreadsheet",
    }
    [structured_entity] = _database_rows(
        select(StructuredKnowledgeEntityModel).where(
            StructuredKnowledgeEntityModel.document_id == document_id,
            StructuredKnowledgeEntityModel.entity_type == "api_endpoint",
        )
    )
    assert json.loads(structured_entity.source_location_json) == source_location

    session_generator = get_database_session()
    db = next(session_generator)
    try:
        retrieved = RetrievalService(db).retrieve(
            user_id=user_id,
            query="OfficeProvenanceAlpha",
        )
        debug_context = retrieved_context_for_graph(retrieved)
    finally:
        session_generator.close()
    assert debug_context[0]["source_location_json"] == source_location

    conversation = client.post("/conversations", json={"title": "Office provenance RAG"})
    assert conversation.status_code == 201
    run = client.post(
        f"/conversations/{conversation.json()['id']}/runs",
        json={"message": "What mentions OfficeProvenanceAlpha?"},
    )

    assert run.status_code == 200
    assert run.json()["citations"]
    assert run.json()["citations"][0]["source_location_json"] == source_location


def test_powerpoint_upload_ingests_slide_source_location_and_cites_it(monkeypatch) -> None:  # noqa: ANN001
    client = _client(monkeypatch)
    _signup_login(client, "pptx-provenance-owner@example.com")
    kb_id = _create_personal_knowledge_base(client, "PowerPoint Provenance KB")
    document = _upload_document(
        client,
        data={"title": "Slide Provenance", "knowledge_base_id": kb_id},
        files={
            "file": (
                "slides.pptx",
                _pptx_bytes(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    assert document.status_code == 201
    payload = document.json()
    assert payload["source_type"] == "presentation"
    assert payload["parser_name"] == "python_pptx_markdown_v1"

    ingest = client.post(f"/documents/{payload['id']}/ingest")
    assert ingest.status_code == 200
    chunks = _database_rows(
        select(DocumentChunkModel)
        .where(DocumentChunkModel.document_id == payload["id"])
        .order_by(DocumentChunkModel.ordinal)
    )
    slide_chunk = next(item for item in chunks if "OfficeSlideAlpha" in item.content)
    source_location = json.loads(slide_chunk.source_location_json)
    assert source_location["source_type"] == "presentation"
    assert source_location["slide_number"] == 1
    assert source_location["slide_title"] == "Office Slide Plan"
    assert source_location["shape_index"] is not None

    conversation = client.post("/conversations", json={"title": "PowerPoint provenance RAG"})
    assert conversation.status_code == 201
    run = client.post(
        f"/conversations/{conversation.json()['id']}/runs",
        json={"message": "What depends on OfficeSlideAlpha?"},
    )

    assert run.status_code == 200
    assert run.json()["citations"]
    assert run.json()["citations"][0]["source_location_json"] == source_location


def test_langgraph_academy_pdf_regression_extracts_real_text_when_available() -> None:
    sample = (
        Path.home() / "Downloads/LangChain_Academy_-_Introduction_to_LangGraph_-_Motivation.pdf"
    )
    if not sample.exists():
        pytest.skip("local LangGraph Academy sample PDF is not available")

    parsed = parse_uploaded_pdf(
        filename=sample.name,
        content_type="application/pdf",
        content=sample.read_bytes(),
    )
    chunks = _chunk_pdf_text(parsed.content)
    entities = {name for chunk, *_ in chunks for name in _extract_entity_names(chunk)}

    assert parsed.parser_name in {"pypdf_text_v2", "pymupdf_text_v1"}
    assert parsed.page_count == 17
    assert "LangChain Academy" in parsed.content
    assert "LangGraph" in parsed.content
    assert len(chunks) >= 10
    assert len(entities) >= 10


def test_pdf_upload_ingest_and_conversation_retrieval_pipeline(monkeypatch) -> None:  # noqa: ANN001
    client = _client(monkeypatch)
    _signup_login(client, "pdf-rag-owner@example.com")
    kb_id = _create_personal_knowledge_base(client, "PDF RAG KB")

    resume_phrase = "Heecheon Park builds FastAPI and LangGraph product systems"
    document = _upload_document(
        client,
        data={"title": "Resume PDF", "knowledge_base_id": kb_id},
        files={
            "file": (
                "resume.pdf",
                _compressed_text_pdf(f"{resume_phrase} with permission-aware retrieval."),
                "application/pdf",
            )
        },
    )
    assert document.status_code == 201

    ingest = client.post(f"/documents/{document.json()['id']}/ingest")
    assert ingest.status_code == 200
    assert ingest.json()["chunk_count"] == 1

    conversation = client.post("/conversations", json={"title": "Resume PDF RAG"})
    assert conversation.status_code == 201
    run = client.post(
        f"/conversations/{conversation.json()['id']}/runs",
        json={"message": "Tell me about me from my uploaded resume."},
    )

    assert run.status_code == 200
    payload = run.json()
    assert payload["citations"]
    assert payload["citations"][0]["document_id"] == document.json()["id"]
    assert payload["citations"][0]["source_filename"] == "resume.pdf"
    assert payload["citations"][0]["source_page"] == 1
    assert payload["citations"][0]["source_location_json"] is None
    assert resume_phrase in payload["citations"][0]["snippet"]


def test_owner_can_delete_uploaded_pdf_and_cleanup_ingestion_artifacts(monkeypatch) -> None:  # noqa: ANN001
    owner = _client(monkeypatch)
    reader = _client(monkeypatch)
    _signup_login(owner, "delete-pdf-owner@example.com")
    reader_id = _signup_login(reader, "delete-pdf-reader@example.com")
    kb_id = _create_personal_knowledge_base(owner, "Delete PDF KB")

    document = _upload_document(
        owner,
        data={"title": "Delete Me PDF", "knowledge_base_id": kb_id},
        files={
            "file": (
                "delete-me.pdf",
                _text_pdf("Delete Cleanup PDF mentions FastAPI and GET /cleanup."),
                "application/pdf",
            )
        },
    )
    assert document.status_code == 201
    document_id = document.json()["id"]

    permission = owner.patch(
        f"/documents/{document_id}/permissions",
        json={"user_id": reader_id, "can_read": True},
    )
    assert permission.status_code == 200
    ingest = owner.post(f"/documents/{document_id}/ingest")
    assert ingest.status_code == 200

    conversation = owner.post("/conversations", json={"title": "Delete cleanup RAG"})
    assert conversation.status_code == 201
    run = owner.post(
        f"/conversations/{conversation.json()['id']}/runs",
        json={"message": "What does Delete Cleanup PDF mention?"},
    )
    assert run.status_code == 200
    assert run.json()["citations"]

    assert _database_rows(select(DocumentModel).where(DocumentModel.id == document_id))
    assert _database_rows(
        select(DocumentChunkModel).where(DocumentChunkModel.document_id == document_id)
    )
    assert _database_rows(
        select(ExtractionRunModel).where(ExtractionRunModel.document_id == document_id)
    )
    assert _database_rows(
        select(EntityMentionModel).where(EntityMentionModel.document_id == document_id)
    )
    assert _database_rows(
        select(EntityRelationshipModel).where(EntityRelationshipModel.document_id == document_id)
    )
    assert _database_rows(
        select(StructuredKnowledgeEntityModel).where(
            StructuredKnowledgeEntityModel.document_id == document_id
        )
    )
    assert _database_rows(
        select(DocumentPermissionModel).where(DocumentPermissionModel.document_id == document_id)
    )
    assert _database_rows(select(CitationModel).where(CitationModel.document_id == document_id))

    deleted = owner.delete(f"/documents/{document_id}")

    assert deleted.status_code == 204
    assert owner.get(f"/documents/{document_id}").status_code == 404
    assert _database_rows(select(DocumentModel).where(DocumentModel.id == document_id)) == []
    assert (
        _database_rows(
            select(DocumentChunkModel).where(DocumentChunkModel.document_id == document_id)
        )
        == []
    )
    assert (
        _database_rows(
            select(ExtractionRunModel).where(ExtractionRunModel.document_id == document_id)
        )
        == []
    )
    assert (
        _database_rows(
            select(EntityMentionModel).where(EntityMentionModel.document_id == document_id)
        )
        == []
    )
    assert (
        _database_rows(
            select(EntityRelationshipModel).where(
                EntityRelationshipModel.document_id == document_id
            )
        )
        == []
    )
    assert (
        _database_rows(
            select(StructuredKnowledgeEntityModel).where(
                StructuredKnowledgeEntityModel.document_id == document_id
            )
        )
        == []
    )
    assert (
        _database_rows(
            select(DocumentPermissionModel).where(
                DocumentPermissionModel.document_id == document_id
            )
        )
        == []
    )
    assert (
        _database_rows(select(CitationModel).where(CitationModel.document_id == document_id)) == []
    )


def test_non_owner_cannot_delete_document_without_manage_authorization(monkeypatch) -> None:  # noqa: ANN001
    owner = _client(monkeypatch)
    reader = _client(monkeypatch)
    _signup_login(owner, "delete-denied-owner@example.com")
    reader_id = _signup_login(reader, "delete-denied-reader@example.com")
    kb_id = _create_personal_knowledge_base(owner, "Delete Denied KB")

    document = _create_document(
        owner,
        json={
            "title": "Private delete guard",
            "content": "reader can see but not delete",
            "knowledge_base_id": kb_id,
        },
    )
    assert document.status_code == 201
    document_id = document.json()["id"]
    grant = owner.patch(
        f"/documents/{document_id}/permissions",
        json={"user_id": reader_id, "can_read": True},
    )
    assert grant.status_code == 200
    assert reader.get(f"/documents/{document_id}").status_code == 200

    denied = reader.delete(f"/documents/{document_id}")

    assert denied.status_code == 404
    assert owner.get(f"/documents/{document_id}").status_code == 200


def test_upload_rejects_unsupported_or_unsafe_input(monkeypatch) -> None:  # noqa: ANN001
    client = _client(monkeypatch)
    _signup_login(client, "pdf-safety@example.com")
    kb_id = _create_personal_knowledge_base(client, "PDF Safety KB")

    unsupported = _upload_document(
        client,
        data={"title": "Docx", "knowledge_base_id": kb_id},
        files={"file": ("notes.docx", b"not supported", "text/plain")},
    )
    assert unsupported.status_code == 415

    legacy_office = _upload_document(
        client,
        data={"title": "Legacy XLS", "knowledge_base_id": kb_id},
        files={"file": ("legacy.xls", b"not supported", "application/vnd.ms-excel")},
    )
    assert legacy_office.status_code == 415

    binary_text = _upload_document(
        client,
        data={"title": "Binary text"},
        files={"file": ("notes.txt", b"hello\x00not text", "text/plain")},
    )
    assert binary_text.status_code == 400

    unsafe_name = _upload_document(
        client,
        data={"title": "Bad name"},
        files={"file": ("../bad.pdf", _text_pdf("Safe text"), "application/pdf")},
    )
    assert unsafe_name.status_code == 400


def test_group_knowledge_base_requires_group_membership(monkeypatch) -> None:  # noqa: ANN001
    owner = _client(monkeypatch)
    member = _client(monkeypatch)
    outsider = _client(monkeypatch)
    _signup_login(owner, "group-kb-owner@example.com")
    member_id = _signup_login(member, "group-kb-member@example.com")
    _signup_login(outsider, "group-kb-outsider@example.com")
    group_id = owner.post("/groups", json={"name": "KB Group"}).json()["id"]
    default_group_kbs = [
        kb
        for kb in owner.get("/knowledge-bases").json()
        if kb["scope"] == "group" and kb["group_id"] == group_id
    ]
    assert len(default_group_kbs) == 1
    assert default_group_kbs[0]["name"] == "KB Group Knowledge"

    assert (
        owner.post(
            f"/groups/{group_id}/members",
            json={"user_id": member_id, "role": "viewer"},
        ).status_code
        == 204
    )

    kb = owner.post(
        "/knowledge-bases",
        json={"name": "Group KB", "scope": "group", "group_id": group_id},
    )
    assert kb.status_code == 201
    assert kb.json()["group_id"] == group_id

    member_denied = member.post(
        "/knowledge-bases",
        json={"name": "Member Denied", "scope": "group", "group_id": group_id},
    )
    assert member_denied.status_code == 403

    denied = outsider.post(
        "/knowledge-bases",
        json={"name": "Denied", "scope": "group", "group_id": group_id},
    )
    assert denied.status_code == 403
    assert outsider.get("/knowledge-bases").json() == []


def test_direct_document_writes_to_group_knowledge_base_are_rejected(monkeypatch) -> None:  # noqa: ANN001
    owner = _client(monkeypatch)
    _signup_login(owner, "group-kb-direct-write-owner@example.com")
    group_id = owner.post("/groups", json={"name": "Direct Write Group"}).json()["id"]
    group_kb = next(
        kb
        for kb in owner.get("/knowledge-bases").json()
        if kb["scope"] == "group" and kb["group_id"] == group_id
    )

    create_response = owner.post(
        f"/knowledge-bases/{group_kb['id']}/documents",
        json={"title": "Direct group doc", "content": "should use publish workflow"},
    )
    upload_response = _upload_document(
        owner,
        data={"title": "Direct Upload", "knowledge_base_id": group_kb["id"]},
        files={"file": ("direct.txt", b"group upload", "text/plain")},
    )

    assert create_response.status_code == 422
    assert create_response.json()["detail"] == (
        "group knowledge bases accept documents through publish approval"
    )
    assert upload_response.status_code == 422
    assert upload_response.json()["detail"] == (
        "group knowledge bases accept documents through publish approval"
    )
